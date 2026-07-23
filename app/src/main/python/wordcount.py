#!/usr/bin/env python3



# -*- coding: utf-8 -*-



"""



wordcount.py  --  统一字数统计技能（合并 PDF / CAD / Excel / PPT 等格式）







把任意受支持格式的文件提取出「可编辑文字」，复制进 Word(.docx)，并按



Word「字数统计」对话框的口径统计：



    · 字数



    · 中文字符和朝鲜语单词   (FarEast 字符：CJK + 朝鲜语 Jamo/音节/兼容 + 全角)



    · 非中文单词             (= 字数 − 中文字符和朝鲜语单词)



    · 字符数(不计空格)







设计要点



--------



* 每个文件先被提取成「文本单元」(items)：字符串=一个段落；("table", 行列表)=一张表。



* 计数与建 Word 文档使用同一份 items，因此「统计结果」==「文档里的文字」，



  等价于你把文字复制到 Word 后看 Word 字数统计（算法已用 Word COM 的「字数=101」



  校验通过，见 Excel 任务）。



* 逐段落 / 逐单元格独立计词（段落、单元格、形状之间都是词分隔），与 Word 一致。



* 支持批量：多个文件 / 目录（可递归）；每个文件生成 <原名>_字数统计.docx；



  可选 BATCH_wordcount_summary.png 汇总图；可选 --send 打包发 QQ 邮箱。







支持的格式



----------



  .pdf                -> pdfminer 提取正文



  .xlsx / .xlsm       -> openpyxl 读取（多工作表，逐单元格）



  .pptx               -> python-pptx 读取（幻灯片文字 + 表格；备注默认不含）



  .ppt                -> PowerPoint COM 读取（旧版二进制格式）



  .docx               -> python-docx 读取（段落 + 表格）



  .doc                -> Word COM 读取（旧版二进制格式，需本机安装 Word）



  .txt                -> 按行分段为段落



  .png/.jpg/.jpeg/.bmp/.tif/.tiff/.gif/.webp -> RapidOCR 图片 OCR 取文字



  .dwg / .dxf         -> LibreDWG(dwg2dxf) 转 DXF + 容错扫描器取矢量文字 +



                         OLE 嵌入对象（Excel/Word/PPT 粘贴进 CAD 的"图片文字"）



  .zip/.rar/.7z/.tar/.tar.gz/.tgz/.tar.bz2/.tar.xz/.gz/.bz2/.xz



                      -> 自动解压到临时目录，递归统计内部所有受支持文件



                         （支持嵌套压缩包；ZIP 中文文件名自动纠正编码）







依赖（隔离 venv，首次自动安装）：pdfminer.six, python-docx, pillow, openpyxl,



python-pptx, ezdxf, olefile, rapidocr_onnxruntime（图片 OCR）, py7zr（.7z 解压）,



rarfile（.rar 解压，需本机有 unrar/bsdtar）。CAD 还需本机有



dwg2dxf.exe（LibreDWG），可用 --converter 指定，或放在本技能/工具目录自动发现。



.doc 与 .ppt 依赖本机安装 Word / PowerPoint 的 COM 自动化；图片 OCR 模型随



rapidocr_onnxruntime 包内置，无需联网下载。



"""



import os
import traceback



import sys


# ═════════════════════════════════════════════════════
# v1.1.26: Android lxml 崩溃拦截器（五层防御终极版）
# ═════════════════════════════════════════════════════
# 根因：lxml 的 C 扩展(.so)在 Android 上触发 fatal 级
# FileNotFoundError(AssetFinder/scripts)，该错误无法被
# Python except 捕获，直接导致进程崩溃。
#
# 历史教训（v1.1.15~v1.1.25 连续11个版本未根治）：
#   v1.1.15~1.18: 只实现 PEP 302 find_module → Python 3.10 跳过
#   v1.1.19:     加了 PEP 451 find_spec 但 raise → import 机制吞异常
#   v1.1.20~21: find_spec 返回假 Spec 但 NameError importlib 未定义
#   v1.1.22~23: 移除 openpyxl+python-pptx（用户依赖路径）
#   v1.1.24~25: create_module 返回空 module + exec_module pass
#                ★ 但仍失败！根因：Chaquopy AssetFinder 通过
#                  sys.path 机制（非 sys.meta_path）加载 lxml，
#                  meta_path 拦截器对 path-based 导入无效！
#
# ★ v1.1.26 五层防御方案 ★
#   第1层: sys.meta_path 拦截器（PEP 302 + PEP 451）— 显式 import
#   第2层: sys.modules 预填充 — 任何 "import lxml" 先查缓存即命中
#   第3层: sys.path 清理 — 删除含 lxml 子目录的 path 条目（阻断 path-based finder）
#   第4层: sys.path_hooks 注入 — 自定义 PathEntryFinder 拦截 lxml 目录
#   第5层: 覆盖 builtins.__import__ — 兜底捕获所有导入尝试
# ═════════════════════════════════════════════════════
import importlib as _importlib
import importlib.util as _iu
from importlib.abc import Loader as _LoaderBase, MetaPathFinder as _MetaPathFinderBase
import types as _types
import builtins as _builtins

# ── 共享工具 ──
_BLOCKED_LOADER = None  # 下面定义后回填

def _is_lxml_name(name):
    """检查模块名是否属于 lxml 家族。"""
    return (name == 'lxml'
            or name.startswith('lxml.')
            or name.startswith('lxml_')
            or name == 'lxml.etree'
            or name in ('lxml.html', 'lxml.sax', 'lxml.isoschematron',
                        'lxml.ElementInclude', 'lxml.cssselect', 'lxml.builder'))


# ══ 第1层：sys.meta_path 拦截器 ══
class _BlockedLoader(_LoaderBase):
    """假 loader：返回空模块，阻止 .so 被加载。"""
    def create_module(self, spec):
        mod = _types.ModuleType(spec.name)
        mod.__loader__ = self
        mod.__package__ = spec.parent
        mod.__spec__ = spec
        return mod  # ← 必须返回非 None！否则 Python 默认机制加载 .so → 崩溃

    def exec_module(self, module):
        pass  # ← 用 pass 不用 raise（raise 会尝试下一个 finder）


_BLOCKED_LOADER = _BlockedLoader()


class _LxmlMetaPathBlocker(_MetaPathFinderBase):
    """第1层：meta_path 拦截所有 lxml 导入。"""

    def find_spec(self, fullname, path, target=None):
        if _is_lxml_name(fullname):
            return _iu.module_spec(fullname, _BLOCKED_LOADER)
        return None

    def find_module(self, fullname, path=None):
        if _is_lxml_name(fullname):
            return self
        return None

    def load_module(self, fullname):
        if fullname in sys.modules:
            return sys.modules[fullname]
        mod = _types.ModuleType(fullname)
        mod.__loader__ = self
        sys.modules[fullname] = mod
        return mod


# 插入到 meta_path 最前面
sys.meta_path.insert(0, _LxmlMetaPathBlocker())


# ══ 第2层：sys.modules 预填充 ══
_dummy_lxml = _types.ModuleType('lxml')
_dummy_lxml.__path__ = []
_dummy_lxml.__loader__ = _BLOCKED_LOADER
sys.modules['lxml'] = _dummy_lxml
for _sub in ('lxml.etree', 'lxml.html', 'lxml.sax', 'lxml.isoschematron',
             'lxml.ElementInclude', 'lxml.cssselect', 'lxml.builder'):
    _m = _types.ModuleType(_sub)
    _m.__loader__ = _BLOCKED_LOADER
    sys.modules[_sub] = _m


# ══ 第3层：sys.path 清理 ══
# Chaquopy 将 pip 包安装在 assets 目录下，通过 sys.path 暴露。
# 如果某个 path 条目下存在 lxml/ 子目录，path-based finder 会直接从那里加载，
# 完全绕过 sys.meta_path！必须将这些条目从 sys.path 中移除。
_cleaned_path = []
for _pentry in list(sys.path):
    # 检查该路径下是否存在 lxml 相关文件/目录
    if _pentry:
        # 尝试检测 lxml 目录或 .so 文件
        _lxml_dir = None
        for _candidate in [
            _pentry + '/lxml',
            _pentry + '/lxml.etree',
        ]:
            try:
                import os as _os
                if _os.path.isdir(_candidate) or _os.path.isfile(_candidate + '.so') or _os.path.isfile(_candidate + '.pyd'):
                    _lxml_dir = _candidate
                    break
            except Exception:
                pass
        if _lxml_dir is None:
            _cleaned_path.append(_pentry)
        else:
            # 发现 lxml 所在路径，整条移除
            import logging as _logging
            _logging.getLogger(__name__).warning(
                "lxml-defender L3: removed sys.path entry containing lxml: %s", _pentry)
sys.path[:] = _cleaned_path
# 注意：不再 del 循环变量（_pentry/_sub/_m 可能在空列表时未赋值 → NameError）


# ══ 第4层：sys.path_hooks 注入 ══
# 即使 lxml 目录残留在 sys.path 中（L3 遗漏），自定义的 PathEntryFinder
# 会在目录级查找时拦截 lxml 子包。
class _LxmlPathEntryFinder:
    """伪装成 PathEntryFinder，当被要求查找 lxml 时返回假 loader。"""

    def find_module(self, fullname, path=None):
        if _is_lxml_name(fullname):
            return _BLOCKED_LOADER
        return None

    def find_spec(self, fullname, target=None):
        if _is_lxml_name(fullname):
            return _iu.module_spec(fullname, _BLOCKED_LOADER)
        return None

    def __repr__(self):
        return '<LxmlPathEntryFinder>'


_LXML_PATH_HOOK_ENTRY = _LxmlPathEntryFinder()


def _lxml_path_hook(path):
    """如果路径看起来像包含 lxml 的包目录，返回拦截 finder。"""
    if not path:
        raise ImportError()
    try:
        # 检查路径是否为 lxml 相关
        _basename = path.rsplit('/', 1)[-1] if '/' in path else path.rsplit('\\', 1)[-1]
        if _basename in ('lxml', 'lxml.etree', 'lxml.html'):
            import logging as _logging
            _logging.getLogger(__name__).warning(
                "lxml-defender L4: intercepted path_hook for lxml: %s", path)
            return _LXML_PATH_HOOK_ENTRY
    except Exception:
        pass
    raise ImportError("lxml-defender: not an lxml path")


# 在 path_hooks 最前面注入我们的钩子
sys.path_hooks.insert(0, _lxml_path_hook)

# 清除 path_importer_cache 中可能的 lxml 缓存
_for_removal = [k for k in sys.path_importer_cache if k and 'lxml' in k]
for _k in _for_removal:
    del sys.path_importer_cache[_k]
# 不再 del 任何变量（v1.1.26~1.1.29 连续3个版本被 del 坑的教训）


# ══ 第5层：覆盖 builtins.__import__ ══
_original_import = _builtins.__import__


def _safe_import(name, globals=None, locals=None, fromlist=(), level=0):
    """最终兜底：拦截所有 lxml 导入尝试。"""
    if _is_lxml_name(name):
        import logging as _logging
        _logging.getLogger(__name__).warning(
            "lxml-defender L5: blocked __import__(%s)", name)
        # 确保已注册到 sys.modules
        if name not in sys.modules:
            _mod = _types.ModuleType(name)
            _mod.__loader__ = _BLOCKED_LOADER
            sys.modules[name] = _mod
        if fromlist:
            # 处理 "from lxml import etree" 形式
            return sys.modules[name]
        return sys.modules[name]
    return _original_import(name, globals, locals, fromlist, level)


_builtins.__import__ = _safe_import

# ═════════════════════════════════════════════════════
# 注意：不再 del 任何变量（v1.1.26~1.1.29 连续3个版本的教训）
#   _p  (v1.1.26): 循环变量未赋值 → NameError
#   _k  (v1.1.28): 循环变量未赋值 → NameError  
#   _original_import (v1.1.29): 被 _safe_import 闭包引用，del 后运行时 NameError
# 这些 _ 前缀私有变量留在模块命名空间无任何副作用，不值得冒险清理


import re



import glob



import json



import argparse



import subprocess



import zipfile



import tarfile



import shutil



import tempfile







# ---- 可移植的运行时路径（不写死用户名）----------------------------------



_HOME = os.path.expanduser("~")



BASE_PY = os.path.join(_HOME, ".workbuddy", "binaries", "python", "versions", "3.13.12", "python.exe")



VENV_DIR = os.path.join(_HOME, ".workbuddy", "binaries", "python", "envs", "default")



VENV_PY = os.path.join(VENV_DIR, "Scripts", "python.exe")



SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))



SCRIPTS_DIR = os.path.join(SKILL_DIR, "scripts")



if SCRIPTS_DIR not in sys.path:



    sys.path.insert(0, SCRIPTS_DIR)











# ---------------------------------------------------------------------------



# 环境自举



# ---------------------------------------------------------------------------



def ensure_env():



    if not os.path.exists(VENV_PY):



        subprocess.run([BASE_PY, "-m", "venv", VENV_DIR], check=True)



    test = ("import pdfminer.high_level, docx, PIL, openpyxl, pptx, ezdxf, olefile, "



            "py7zr, fitz; "



            "print('deps ok')")



    r = subprocess.run([VENV_PY, "-c", test], capture_output=True, text=True)



    if r.returncode != 0:



        log = os.path.join(SCRIPTS_DIR, "wordcount_pip_install.log")



        subprocess.run(



            [VENV_PY, "-m", "pip", "install", "--timeout", "180",



             "-i", "https://pypi.tuna.tsinghua.edu.cn/simple",



             "pdfminer.six", "python-docx", "pillow", "openpyxl",



             "python-pptx", "ezdxf", "olefile", "py7zr", "rarfile", "pymupdf"],



            stdout=open(log, "w"), stderr=subprocess.STDOUT, check=True)











def reexec_with_venv():



    # v1.3.15：打包成 PyInstaller exe 后 sys.frozen 为真，依赖已全部内嵌，
    # 不能再尝试用 venv 的 python 重新执行自身（临时目录里没有 wordcount.py）
    if getattr(sys, "frozen", False):
        return
    if os.path.abspath(sys.executable).lower() != os.path.abspath(VENV_PY).lower():



        ensure_env()



        os.execv(VENV_PY, [VENV_PY, os.path.abspath(__file__)] + sys.argv[1:])











# ---------------------------------------------------------------------------



# Word 口径计数（已用 Word COM 字数=101 校验）



# ---------------------------------------------------------------------------



_FAR = (



    r"\u1100-\u11FF"   # Hangul Jamo



    r"\u3000-\u303F"   # CJK 符号与标点



    r"\u3130-\u318F"   # Hangul 兼容 Jamo



    r"\u3400-\u4DBF"   # CJK 扩展 A



    r"\u4E00-\u9FFF"   # CJK 基本平面



    r"\uA960-\uA97C"   # Hangul Jamo 扩展 A



    r"\uAC00-\uD7A3"   # Hangul 音节



    r"\uD7B0-\uD7FF"   # Hangul Jamo 扩展 B



    r"\uF900-\uFAFF"   # CJK 兼容



    r"\uFF00-\uFFEF"   # 全角字符



)



FAR_EAST = re.compile("[" + _FAR + "]")



NON_CJK = re.compile("[^" + r"\s" + _FAR + "]+")











def count_unit(s):



    """Return (fe, nc, chars) for a single text unit (paragraph or cell)."""



    fe = len(FAR_EAST.findall(s))



    nc = len(NON_CJK.findall(s))



    ch = len(re.sub(r"\s", "", s))



    return fe, nc, ch











def count_items(items):



    """items: list where each element is either a str (paragraph) or



    ("table", rows) where rows is list[list[str]].







    逐段落 / 逐单元格独立计词 -> 段落、单元格、形状之间的边界都是词分隔，与 Word 一致。



    """



    total_fe = total_nc = total_ch = 0



    for it in items:



        if isinstance(it, str):



            fe, nc, ch = count_unit(it)



            total_fe += fe



            total_nc += nc



            total_ch += ch



        else:  # ("table", rows)



            rows = it[1] if len(it) > 1 else []



            for row in rows:



                for cell in row:



                    fe, nc, ch = count_unit(cell if cell is not None else "")



                    total_fe += fe



                    total_nc += nc



                    total_ch += ch



    return {"fe": total_fe, "nc": total_nc, "chars": total_ch,



            "words": total_fe + total_nc}











# ---------------------------------------------------------------------------



# 建 Word 文档（与计数用同一份 items，保证统计==文档文字）



# ---------------------------------------------------------------------------



def _append_items(doc, items):



    """把 items（段落/表格混合列表）写入 doc。供 build_docx / build_merged_docx 复用。"""



    for it in items:



        if isinstance(it, str):



            doc.add_paragraph(it if it != "" else "")



        else:  # table



            rows = it[1] if len(it) > 1 else []



            if not rows:



                continue



            nrows = len(rows)



            ncols = max((len(r) for r in rows), default=0)



            if ncols == 0:



                continue



            table = doc.add_table(rows=nrows, cols=ncols)



            table.style = "Table Grid"



            for i, row in enumerate(rows):



                for j, cell in enumerate(row):



                    table.cell(i, j).text = cell if cell is not None else ""



                    if i == 0:



                        for p in table.cell(i, j).paragraphs:



                            for run in p.runs:



                                run.bold = True



            doc.add_paragraph()











def build_docx(title, items, out_path, stats=None, meta=None):



    from docx import Document



    from docx.shared import Pt, RGBColor



    doc = Document()



    doc.add_heading(title, level=1)



    if meta and meta.get("img_items"):



        # 文字层段落 + 一段非计词的标题（说明下方是图片页 OCR 文字），再接 OCR 段落



        _append_items(doc, meta.get("text_items") or items)



        n = len(meta.get("img_pages", []))



        rng = meta.get("img_pages_ranges", "")



        doc.add_heading(



            "（以下为图片页 OCR 识别文字：第 %s 页，共 %d 页）" % (rng, n),



            level=2)



        _append_items(doc, meta["img_items"])



    else:



        _append_items(doc, items)







    if stats is not None:



        doc.add_heading("Word 字数统计结果（与 Word 字数统计对话框一致）", level=2)



        p = doc.add_paragraph()



        p.add_run("字数：").bold = True



        p.add_run(f"{stats['words']}\n")



        p.add_run("中文字符和朝鲜语单词：").bold = True



        p.add_run(f"{stats['fe']}\n")



        p.add_run("非中文单词：").bold = True



        p.add_run(f"{stats['nc']}\n")



        p.add_run("字符数（不计空格）：").bold = True



        p.add_run(f"{stats['chars']}")



        note = doc.add_paragraph()



        nr = note.add_run(



            "说明：以上对应本文件提取出的可编辑文字；"



            "『字数 = 中文字符和朝鲜语单词 + 非中文单词』。"



            "表内数字按非中文单词计入，中文及全角标点按中文字符计入。")



        nr.italic = True



        nr.font.size = Pt(9)



        nr.font.color.rgb = RGBColor(0x66, 0x66, 0x66)







    doc.save(out_path)











def build_merged_docx(results, out_path, png_path=None):



    """把多个文件的统计结果合并成一个 Word 文档：



       * 每个文件：标题(文件名) + 提取出的文字/表格 + 该文件字数统计



       * 末尾：汇总表(各文件 + 合计) + 数字统计汇总图（如有，置于文档最后）



    results: list of dict{src, items, stats, meta}



    """



    from docx import Document



    from docx.shared import Pt, RGBColor, Inches



    from docx.enum.text import WD_ALIGN_PARAGRAPH







    doc = Document()



    doc.add_heading("WordCount 字数统计汇总报告", level=0)







    for r in results:



        src = r.get("src", "")



        items = r.get("items", []) or []



        stats = r.get("stats", {})



        rmeta = r.get("meta", {}) or {}



        doc.add_heading(os.path.basename(src), level=1)



        if rmeta and rmeta.get("img_items"):



            _append_items(doc, rmeta.get("text_items") or items)



            n = len(rmeta.get("img_pages", []))



            rng = rmeta.get("img_pages_ranges", "")



            doc.add_heading(



                "（以下为图片页 OCR 识别文字：第 %s 页，共 %d 页）" % (rng, n),



                level=2)



            _append_items(doc, rmeta["img_items"])



        else:



            _append_items(doc, items)



        # 该文件统计



        p = doc.add_paragraph()



        p.add_run("字数：").bold = True



        p.add_run(f"{stats.get('words', 0)}    ")



        p.add_run("中文字符和朝鲜语单词：").bold = True



        p.add_run(f"{stats.get('fe', 0)}    ")



        p.add_run("非中文单词：").bold = True



        p.add_run(f"{stats.get('nc', 0)}    ")



        p.add_run("字符数(不计空格)：").bold = True



        p.add_run(f"{stats.get('chars', 0)}")



        doc.add_paragraph()







    # 汇总表



    doc.add_heading("字数统计汇总", level=1)



    cols = ["文件", "字数", "中文字符和朝鲜语单词", "非中文单词", "字符数(不计空格)"]



    rows = [[os.path.basename(r.get("src", ""))] +



            [str(r.get("stats", {}).get(k, 0)) for k in ("words", "fe", "nc", "chars")]



            for r in results]



    tot = {"words": 0, "fe": 0, "nc": 0, "chars": 0}



    for r in results:



        for k in tot:



            tot[k] += r.get("stats", {}).get(k, 0)



    rows.append(["合计"] + [str(tot[k]) for k in ("words", "fe", "nc", "chars")])







    table = doc.add_table(rows=len(rows), cols=len(cols))



    table.style = "Table Grid"



    for j, c in enumerate(cols):



        table.cell(0, j).text = c



        for p in table.cell(0, j).paragraphs:



            for run in p.runs:



                run.bold = True



    for i, row in enumerate(rows[1:], start=1):



        for j, val in enumerate(row):



            table.cell(i, j).text = val



    # 合计行加粗



    last = len(rows) - 1



    for j in range(len(cols)):



        for p in table.cell(last, j).paragraphs:



            for run in p.runs:



                run.bold = True







    # 数字统计图放在文档最后



    if png_path and os.path.exists(png_path):



        doc.add_paragraph()



        try:



            from PIL import Image



            with Image.open(png_path) as im:



                iw, ih = im.size



            # 按页宽缩放（A4 可用宽度约 6.3 英寸）



            max_w = Inches(6.3)



            doc.add_picture(png_path, width=max_w)



            doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER



            cap = doc.add_paragraph("图：各文件字数统计汇总")



            cap.alignment = WD_ALIGN_PARAGRAPH.CENTER



            for run in cap.runs:



                run.font.size = Pt(9)



                run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)



        except Exception:



            pass







    doc.save(out_path)











# ---------------------------------------------------------------------------



# 汇总 PNG（Pillow，含 CJK 字体回退）



# ---------------------------------------------------------------------------



def get_cjk_font(size):



    from PIL import ImageFont



    candidates = [



        "C:/Windows/Fonts/msyh.ttc",



        "C:/Windows/Fonts/simhei.ttf",



        "C:/Windows/Fonts/simsun.ttc",



        "C:/Windows/Fonts/simfang.ttf",



        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",



    ]



    for c in candidates:



        if os.path.exists(c):



            try:



                return ImageFont.truetype(c, size)



            except Exception:



                pass



    try:



        return ImageFont.truetype("arial.ttf", size)



    except Exception:



        return ImageFont.load_default()











def render_summary_png(results, out_path):



    from PIL import Image, ImageDraw, ImageFont



    font = get_cjk_font(20)



    cols = ["文件", "字数", "中文字符和朝鲜语单词", "非中文单词", "字符数(不计空格)"]



    rows_data = [[os.path.basename(r["src"])] +



                 [str(r["stats"][k]) for k in ("words", "fe", "nc", "chars")]



                 for r in results]



    tot = {"words": 0, "fe": 0, "nc": 0, "chars": 0}



    for r in results:



        for k in tot:



            tot[k] += r["stats"][k]



    rows_data.append(["合计"] + [str(tot[k]) for k in ("words", "fe", "nc", "chars")])







    pad = 16



    line_h = 34



    col_w = []



    for ci, c in enumerate(cols):



        w = font.getbbox(c)[2]



        for row in rows_data:



            w = max(w, font.getbbox(row[ci])[2])



        col_w.append(w + 24)



    row_w = sum(col_w) + pad * 2



    row_h = (len(rows_data) + 1) * line_h + pad * 2



    img = Image.new("RGB", (row_w, row_h), (255, 255, 255))



    d = ImageDraw.Draw(img)



    x = pad



    y = pad



    for ci, c in enumerate(cols):



        d.text((x, y), c, fill=(0, 0, 0), font=font)



        x += col_w[ci]



    d.line([(pad, y + line_h - 6), (row_w - pad, y + line_h - 6)], fill=(200, 200, 200), width=1)



    y += line_h



    for row in rows_data:



        x = pad



        for ci, val in enumerate(row):



            d.text((x, y), val, fill=(0, 0, 0), font=font)



            x += col_w[ci]



        y += line_h



    img.save(out_path)











# ---------------------------------------------------------------------------



# 提取器：PDF



# ---------------------------------------------------------------------------



def _format_page_ranges(pages):



    """把 1-based 页码列表压成『1-3、5、7-9』这样的中文区间串。







    连续的页码合并为『X-Y』；不连续的单独列出，用顿号『、』分隔。



    """



    pages = sorted(set(int(p) for p in pages))



    if not pages:



        return ""



    ranges = []



    start = prev = pages[0]



    for p in pages[1:]:



        if p == prev + 1:



            prev = p



        else:



            ranges.append((start, prev))



            start = prev = p



    ranges.append((start, prev))



    parts = [str(s) if s == e else "%d-%d" % (s, e) for s, e in ranges]



    return "、".join(parts)











def _render_pdf_pages_base64(path):
    """用 PyMuPDF(fitz) 将 PDF 每页渲染为 PNG 图，返回 base64 编码列表。

    专供 Android Kotlin 端调用（PdfOcrEngine 路径4）：
      Kotlin 无法用系统 PdfRenderer / PdfiumAndroid 渲染时
      → 调此函数获取页面图 → ML Kit OCR 识别文字

    返回 JSON 字符串: {"ok": bool, "pages": int, "images": ["base64...", ...], "error": str|None}
    若 fitz 不可用或打开失败 → ok=false + error 原因
    """
    import base64, json
    result = {"ok": False, "pages": 0, "images": [], "error": None}
    try:
        import fitz
    except ImportError:
        result["error"] = "pymupdf_not_installed"
        return json.dumps(result)

    try:
        doc = fitz.open(path)
    except Exception as e:
        result["error"] = "fitz_open_failed: " + str(e)[:200]
        return json.dumps(result)

    try:
        page_count = doc.page_count
        result["pages"] = page_count
        mat = fitz.Matrix(2.0, 2.0)  # 2x 缩放足够 OCR 识别
        images = []
        for i in range(page_count):
            try:
                page = doc[i]
                pix = page.get_pixmap(matrix=mat)
                png_bytes = pix.tobytes(output="png")
                if len(png_bytes) > 1000:  # 太小的可能是空白页
                    b64 = base64.b64encode(png_bytes).decode("ascii")
                    images.append(b64)
            except Exception:
                continue
        result["images"] = images
        if images:
            result["ok"] = True
        else:
            result["error"] = "no_pages_rendered"
    except Exception as e:
        result["error"] = "render_error: " + str(e)[:200]
    finally:
        try:
            doc.close()
        except Exception:
            pass
    return json.dumps(result)


def _pdf_image_ocr(path):



    """扫描 PDF，找出『整页都是图片（无可提取文字）』的页，逐页渲染后用



    RapidOCR 识别文字。







    返回 (img_pages, img_items)：



      img_pages : 这些页的 1-based 页码列表（仅当该页 OCR 出非空文字才计入）



      img_items : 这些页 OCR 出来的文字段落（与 extract_image 同口径，



                  每识别出一段文字记一个段落，交给统一计数逻辑）



    若本机未安装 PyMuPDF(fitz) 或打开失败，返回 ([], []) 优雅降级为纯文字层提取。



    """



    try:



        import fitz



    except Exception:



        return [], []



    try:



        doc = fitz.open(path)



    except Exception:



        return [], []



    img_pages = []



    img_items = []



    # 2.5x 缩放提升小字号识别率；仅渲染无可提取文字的页



    mat = fitz.Matrix(2.5, 2.5)



    tmpdir = tempfile.mkdtemp(prefix="wc_pdfimg_")



    try:



        for i in range(doc.page_count):



            page = doc[i]



            page_text = (page.get_text("text") or "").strip()



            if page_text:



                # 该页本身有可提取文字 -> 视为文字页，不 OCR（避免重复计词）



                continue



            try:



                pix = page.get_pixmap(matrix=mat)



                png = os.path.join(tmpdir, "p%d.png" % (i + 1))



                pix.save(png)



                items = extract_image(png)



            except Exception:



                items = []



            if items:



                img_pages.append(i + 1)



                img_items.extend(items)



    finally:



        try:



            doc.close()



        except Exception:



            pass



        shutil.rmtree(tmpdir, ignore_errors=True)



    return img_pages, img_items











def extract_pdf(path):



    """从 PDF 提取可统计文字。







    规则：



    * 文字层（pdfminer 提取的正文）照常计入；



    * 若某页『整页都是图片、没有可提取文字』（扫描件/纯图片页），则渲染该页用



      RapidOCR 图片转文字后再计入（与统计普通图片的逻辑一致）；



    * 合并后统计 == 文字层文字 + 图片页 OCR 文字，等价于你把全部文字复制到



      Word 后看 Word 字数统计。







    返回 (items, meta)：



      items : 文字层段落 + 图片页 OCR 段落（合并，供计数与建 Word）



      meta  : 含图片页明细，供 GUI 展开查看：



              {"img_pages":[..], "img_pages_ranges":"..", "img_stats":{...},



               "img_items":[..], "text_items":[..]}



              没有图片页时 meta 为空 dict。



    """



    from pdfminer.high_level import extract_text



    try:



        text = extract_text(path)



    except Exception:



        text = ""



    text_items = []



    for para in re.split(r"\n\s*\n", text):



        c = para.strip()



        if c:



            text_items.append(c)







    img_pages, img_items = _pdf_image_ocr(path)







    items = text_items + img_items



    meta = {}



    if img_pages:



        meta = {



            "img_pages": img_pages,



            "img_pages_ranges": _format_page_ranges(img_pages),



            "img_stats": count_items(img_items),



            "img_items": img_items,



            "text_items": text_items,



        }



        # 图片型 PDF：没有任何可提取文字层，全部靠 OCR。



        # 工程图纸 OCR（尤其中文）识别率较低，统计可能不全。



        if not text_items:



            meta["image_only"] = True



            meta["image_only_note"] = (



                "此 PDF 为图片型（不含可提取的文字层），文字完全经 OCR 识别。"



                "工程图中文字（尤其中文）OCR 识别率较低，统计可能不全；"



                "建议提供“文字型”PDF（导出时保留文字）以获得准确统计。"



            )



    return items, meta











# ---------------------------------------------------------------------------



# 提取器：Adobe Illustrator (.ai)



# ---------------------------------------------------------------------------



def extract_ai(path):



    """从 Adobe Illustrator (.ai) 提取文字。







    现代 .ai（默认保存）会内嵌 PDF 兼容数据，PyMuPDF(fitz) 可直接读取其文字层；



    旧版 / 另存为『不含 PDF 内容』的 .ai 仅含 EPS/PS，fitz 读到的只有一句



    “此文件保存时不含 PDF 内容”的说明，此时回退到 pdfminer 试探解析。







    返回 (items, meta)：



      meta["pages"] = 1（设计图为单页）



      meta["ai_no_text"] = True 时表示未嵌入可提取文字层（旧版/另存为无 PDF）



    """



    def _split_text(blob):



        out = []



        for para in re.split(r"\n\s*\n", blob):



            c = para.strip()



            if c:



                out.append(c)



        return out







    # 常见的“无 PDF 内容”说明句前缀（多语言），需剔除，不算作真实文字



    _no_pdf_markers = (



        "toto je soubor", "this file was saved", "this is an adobe",



        "this file does not contain", "il file è stato salvato",



        "este archivo se guardó", "ce fichier a été enregistré",



        "diese datei wurde gespeichert",



    )







    text_items = []



    meta = {"pages": 1, "pages_reason": None}



    try:



        import fitz



        doc = fitz.open(path)



        for page in doc:



            text_items += _split_text(page.get_text())



    except Exception:



        text_items = []



    # 自动识别未拿到文字时，强制按 PDF 解析（兼容“内嵌 PDF 的 .ai”或 PDF 改名 .ai）



    if not text_items:



        try:



            import fitz



            doc = fitz.open(path, filetype="pdf")



            for page in doc:



                text_items += _split_text(page.get_text())



        except Exception:



            pass







    if not text_items:



        # 回退：尝试把 .ai 当作 PDF 用 pdfminer 解析（部分内嵌 PDF 的 .ai）



        try:



            from pdfminer.high_level import extract_text



            t = extract_text(path)



            text_items = _split_text(t)



        except Exception:



            text_items = []







    # 剔除“无 PDF 内容”一类说明句（fitz / pdfminer 均可能读到），保留真实文字



    real_items = []



    for s in text_items:



        low = s.strip().lower()



        if any(mk in low for mk in _no_pdf_markers):



            continue



        real_items.append(s)



    text_items = real_items







    if not text_items:



        meta["ai_no_text"] = True



        meta["ai_note"] = (



            "该 .ai 文件未嵌入可提取的文字层（旧版 Illustrator / 另存为不含 PDF）。"



            "请改用“包含 PDF 兼容内容”的方式重新保存（Illustrator 中另存为 .ai 时"



            "确保勾选『包含 PDF』/ 使用较新版本格式），或用 fitz 可识别的版本后重试。"



        )



    return text_items, meta











# ---------------------------------------------------------------------------



# 提取器：CorelDRAW (.cdr)



# ---------------------------------------------------------------------------



def _extract_text_from_xml(data):



    """从 XML 字节里收集所有可见文本（元素 text / tail）。失败回退正则。"""



    import xml.etree.ElementTree as ET



    parts = []



    try:



        root = ET.fromstring(data)



        for elem in root.iter():



            for t in (elem.text, elem.tail):



                if t:



                    s = t.strip()



                    if s:



                        parts.append(s)



    except Exception:



        try:



            txt = data.decode("utf-8", "ignore")



            parts = [m.strip() for m in re.findall(r">([^<>]+)<", txt) if m.strip()]



        except Exception:



            parts = []



    return parts











def extract_cdr(path):



    """从 CorelDRAW (.cdr) 提取文字。







    CorelDRAW X4(14) 及以后版本的 .cdr 本质是 ZIP 容器，根目录含 content.xml



    （ODF 风格的 XML，文本存于 <text:p>/<text:span> 之类节点），可直接抽取 Unicode



    文字。更早版本为二进制/OLE，本机无 LibreOffice 时无法直接转 PDF，给出原因。







    返回 (items, meta)：



      meta["pages"] = 1（设计图为单页）



      meta["cdr_no_text"] = True / meta["cdr_error"] 时表示无法提取



    """



    import zipfile



    items = []



    meta = {"pages": 1, "pages_reason": None}



    try:



        if zipfile.is_zipfile(path):



            with zipfile.ZipFile(path) as z:



                names = z.namelist()



                xmls = ["content.xml"] if "content.xml" in names else  [n for n in names if n.lower().endswith(".xml")]



                for xn in xmls:



                    try:



                        parts = _extract_text_from_xml(z.read(xn))



                        items += [p for p in parts if p]



                    except Exception:



                        continue



        else:



            meta["cdr_error"] = (



                "该 .cdr 不是 ZIP 容器（可能为旧版二进制/OLE 格式）。"



                "本机未安装 LibreOffice，无法直接转换为可提取文字的 PDF；"



                "建议用 CorelDRAW 另存为『文字型 PDF』后统计。"



            )



    except Exception as e:



        meta["cdr_error"] = "解析 .cdr 失败：%s" % e







    if not items and "cdr_error" not in meta:



        meta["cdr_no_text"] = True



        meta["cdr_note"] = (



            "未能从 .cdr 提取到文字。现代 .cdr(X4+) 通常含 content.xml 可提取；"



            "若文件加密/损坏或版本不兼容，请导出为含文字的 PDF 后统计。"



        )



    return items, meta











# ---------------------------------------------------------------------------



# 提取器：Excel



# ---------------------------------------------------------------------------



def _cell_text(v):
    """将单元格值转为文本。v1.3.15：Excel 日期序列号（int/float，范围约 20000~60000）
    转为中文短日期格式（如「7月1日」），与复制到 Word 后的显示文本一致，
    解决日期列被当纯数字导致 CJK/非CJK 单词大幅偏少的问题。"""
    from datetime import datetime, timedelta

    if v is None:
        return ""
    if isinstance(v, bool):
        return str(v)
    # ── Excel 日期序列号 → 中文短日期 ──
    vi = int(v) if isinstance(v, float) else v
    if isinstance(v, (int, float)) and 20000 < vi < 60000:
        try:
            dt = datetime(1899, 12, 30) + timedelta(days=vi)
            return dt.strftime("%m月%d日")
        except Exception:
            pass
    if isinstance(v, float):
        return str(int(v)) if v.is_integer() else str(v)
    return str(v).strip()
def extract_excel(path, sheet_filter="all"):



    import openpyxl



    wb = openpyxl.load_workbook(path, data_only=True)



    names = wb.sheetnames



    if sheet_filter == "first":



        targets = names[:1]



    elif sheet_filter == "all":

        targets = [s for s in names if wb[s].sheet_state != "hidden"]

    else:

        targets = [s for s in names if s == sheet_filter and wb[s].sheet_state != "hidden"] or [s for s in names if wb[s].sheet_state != "hidden"]


    items = []



    sheet_stats = []



    for name in targets:



        ws = wb[name]



        if ws.max_row == 0 or ws.max_column == 0:



            stats = {"fe": 0, "nc": 0, "chars": 0, "words": 0}



            sheet_stats.append({"name": name, "stats": stats})



            continue



        rows = [[_cell_text(ws.cell(row=r, column=c).value)



                 for c in range(1, ws.max_column + 1)]



                for r in range(1, ws.max_row + 1)]



        total_fe = total_nc = total_ch = 0



        for row in rows:



            for cell in row:



                fe, nc, ch = count_unit(cell)



                total_fe += fe



                total_nc += nc



                total_ch += ch



        stats = {



            "fe": total_fe,



            "nc": total_nc,



            "chars": total_ch,



            "words": total_fe + total_nc,



        }



        sheet_stats.append({"name": name, "stats": stats})



        if stats["words"] > 0:



            items.append(("table", rows))



    return items, sheet_stats











# ---------------------------------------------------------------------------



# 提取器：PPT / PPTX



# ---------------------------------------------------------------------------



def extract_pptx(path, with_notes=False):



    from pptx import Presentation



    prs = Presentation(path)



    items = []







    def _paragraphs(text_frame):



        for p in text_frame.paragraphs:



            txt = (p.text or "").strip()



            if txt:



                items.append(txt)







    for slide in prs.slides:



        for shape in slide.shapes:



            if shape.has_text_frame:



                _paragraphs(shape.text_frame)



            if shape.has_table:



                tbl = shape.table



                rows = [[ (cell.text if cell.text is not None else "")



                          for cell in row.cells ]



                        for row in tbl.rows]



                items.append(("table", rows))



        if with_notes and slide.has_notes_slide:



            ns = slide.notes_slide



            if ns.notes_text_frame is not None:



                _paragraphs(ns.notes_text_frame)



    return items











# ---------------------------------------------------------------------------



# 提取器：CAD (.dwg/.dxf) —— 移植自 cad-wordcount / cad_ole_ocr



# ---------------------------------------------------------------------------



def find_converter(explicit=None):



    if explicit and os.path.exists(explicit):



        return explicit



    name = "dwg2dxf.exe" if sys.platform.startswith("win") else "dwg2dxf"



    roots = [



        os.path.dirname(os.path.abspath(__file__)),



        os.path.join(SKILL_DIR, "..", ".."),



        r"D:\WORKBUDDYWORK\Claw\libredwg",



        os.path.join(_HOME, ".workbuddy", "tools"),



    ]



    for root in roots:



        if not os.path.isdir(root):



            continue



        for dirpath, _, files in os.walk(root):



            if name in files:



                return os.path.join(dirpath, name)



    return None











def dwg_to_dxf(dwg_path, dxf_path, converter):



    cmd = [converter, "-o", dxf_path, dwg_path]



    r = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=240)



    if not os.path.exists(dxf_path):



        alt = os.path.splitext(dwg_path)[0] + ".dxf"



        if os.path.exists(alt) and alt != dxf_path:



            os.replace(alt, dxf_path)



    if not os.path.exists(dxf_path):



        raise RuntimeError("dwg2dxf failed (rc=%s)" % r.returncode)



    return dxf_path











BINARY_GROUPS = {b"310", b"311", b"312", b"313", b"314", b"315",



                 b"316", b"317", b"318", b"319"}











def sanitize_dxf(dxf_path):



    out = dxf_path + "._sanitized.dxf"



    try:



        with open(dxf_path, "rb") as f:



            raw = f.read()



    except Exception:



        return dxf_path



    lines = raw.split(b"\n")



    out_lines = []



    i, n = 0, len(lines)



    while i < n:



        ln = lines[i]



        if ln.strip() in BINARY_GROUPS:



            out_lines.append(ln)



            if i + 1 < n:



                out_lines.append(b"00")



                i += 2



                continue



        out_lines.append(ln)



        i += 1



    with open(out, "wb") as f:



        f.write(b"\n".join(out_lines))



    return out











def clean_mtext(s):



    s = s.replace("\\P", "\n").replace("\\p", "\n")



    s = s.replace("\\~", " ").replace("\\^I", " ").replace("\\^J", " ")



    s = re.sub(r"\\[A-Za-z][^;{}]*;", "", s)



    s = s.replace("\\\\", "\\").replace("{", "").replace("}", "").replace("\\", "")



    return s











def extract_text_custom(dxf_path):



    san = sanitize_dxf(dxf_path)



    try:



        with open(san, "rb") as f:



            raw = f.read()



        try:



            text = raw.decode("utf-8")



        except UnicodeDecodeError:



            text = raw.decode("gb18030", errors="replace")



    except Exception:



        return ""



    lines = text.split("\n")



    entities = []



    cur = None



    i, n = 0, len(lines)



    while i < n - 1:



        code = lines[i].strip()



        val = lines[i + 1]



        i += 2



        if code == "0":



            cur = {"_t": val.strip(), "_g": {}}



            entities.append(cur)



        elif cur is not None:



            cur["_g"].setdefault(code, []).append(val)



    collected = []



    for e in entities:



        t, g = e["_t"], e["_g"]



        if t in ("TEXT", "ATTDEF"):



            vs = g.get("1")



            if vs and vs[0].strip():



                collected.append(vs[0].strip())



        elif t == "MTEXT":



            s = "".join(g.get("1", []) + g.get("3", []))



            s = clean_mtext(s)



            if s.strip():



                collected.append(s.strip())



        elif t == "MULTILEADER":



            s = "".join(g.get("304", []) + g.get("302", []))



            if s.strip():



                collected.append(s.strip())



    seen = set()



    out = []



    for c in collected:



        if c not in seen:



            seen.add(c)



            out.append(c)



    return "\n".join(out)











def extract_text_from_dxf(dxf_path, converter=None):



    try:



        text = extract_text_custom(dxf_path)



        if text.strip():



            return text



    except Exception:



        pass



    import ezdxf



    from ezdxf import DXFStructureError







    def _read(path):



        return ezdxf.readfile(path)







    doc = None



    try:



        doc = _read(dxf_path)



    except DXFStructureError:



        san = sanitize_dxf(dxf_path)



        try:



            doc = _read(san)



        except DXFStructureError:



            try:



                from ezdxf.recover import readfile as recover_readfile



                doc, _auditor = recover_readfile(san)



            except Exception:



                doc = None



    if doc is None and converter and os.path.exists(converter):



        alt = os.path.splitext(dxf_path)[0] + "._min.dxf"



        try:



            subprocess.run([converter, "-m", "-y", "-o", alt, dxf_path],



                           stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=180)



            if os.path.exists(alt):



                doc = ezdxf.readfile(alt)



        except Exception:



            doc = None



    if doc is None:



        return ""



    blocks = doc.blocks



    collected = []







    def grab(entity):



        t = entity.dxftype()



        try:



            if t in ("TEXT", "ATTDEF"):



                txt = entity.dxf.text



            elif t == "MTEXT":



                txt = entity.text if hasattr(entity, "text") else entity.plain_text()



            elif t == "MULTILEADER":



                try:



                    txt = entity.plain_text()



                except Exception:



                    txt = getattr(entity, "text", "")



            else:



                return



        except Exception:



            return



        if txt and txt.strip():



            collected.append(txt.strip())







    for e in doc.modelspace():



        grab(e)



    try:



        for layout in doc.layouts:



            for e in layout:



                grab(e)



    except Exception:



        pass



    try:



        for bname in blocks:



            for e in blocks[bname]:



                grab(e)



    except Exception:



        pass



    seen = set()



    out = []



    for c in collected:



        if c not in seen:



            seen.add(c)



            out.append(c)



    return "\n".join(out)











def _detect_frame_rectangles(doc):



    """从模型空间/图纸空间的『长直线段』重建图框矩形，统计张数。







    适用场景：图框是用 LINE 几何画线构成的 CAD（没有 INSERT 标题块、



    或 dwg2dxf 转换后图纸空间内容丢失、标题块落到模型空间的情况）。



    思路：



      * 收集所有较长的水平/竖直 LINE，按坐标量化后分组；



      * 用水平线段(上下边) + 竖直线段(左右边)两两配对，找出闭合矩形；



      * 仅保留"极大矩形"(不被其它矩形包含的)，排除框内小框(如标题栏)；



      * 若一个超大矩形包含 >=3 个小矩形，则小矩形才是真正的图框(整卷外框)；



      * 最后过滤出"像图纸"的矩形(足够大、宽高不极端)，其数量即页数。



    返回 int（>=0）。



    """



    try:



        from collections import defaultdict



    except Exception:



        return 0



    SNAP = 2







    def sn(v):



        try:



            return int(round(float(v) / SNAP))



        except Exception:



            return 0







    Hy = defaultdict(list)   # y -> [(xmin, xmax), ...]  水平长线



    Vx = defaultdict(list)   # x -> [(ymin, ymax), ...]  竖直长线







    def collect(entities):



        for e in entities:



            if e.dxftype() != "LINE":



                continue



            try:



                x1 = e.dxf.start.x; y1 = e.dxf.start.y



                x2 = e.dxf.end.x;   y2 = e.dxf.end.y



            except Exception:



                continue



            if abs(y1 - y2) < 1e-3:            # 水平线



                L = abs(x2 - x1)



                if L >= 150:



                    Hy[sn(y1)].append((sn(min(x1, x2)), sn(max(x1, x2))))



            elif abs(x1 - x2) < 1e-3:          # 竖直线



                L = abs(y2 - y1)



                if L >= 150:



                    Vx[sn(x1)].append((sn(min(y1, y2)), sn(max(y1, y2))))







    try:



        collect(doc.modelspace())



    except Exception:



        pass



    try:



        for layout in doc.layouts:



            collect(layout)



    except Exception:



        pass



    try:



        for bname in ("*Paper_Space", "*Paper_Space0",



                      "*Paper_Space1", "*Paper_Space2"):



            try:



                collect(doc.blocks[bname])



            except Exception:



                pass



    except Exception:



        pass







    if not Hy or not Vx:



        return 0







    def merge(iv):



        iv.sort()



        out = []



        for a, b in iv:



            if out and a <= out[-1][1] + 1:



                out[-1] = (out[-1][0], max(out[-1][1], b))



            else:



                out.append((a, b))



        return out







    for y in Hy:



        Hy[y] = merge(Hy[y])



    for x in Vx:



        Vx[x] = merge(Vx[x])







    def v_span(x, ylo, yhi):



        for a, b in Vx.get(x, []):



            if a - 1 <= ylo and b + 1 >= yhi:



                return True



        return False







    rects = []



    ylist = sorted(Hy)



    n = len(ylist)



    for i in range(n):



        ya = ylist[i]



        for j in range(i + 1, n):



            yb = ylist[j]



            if yb - ya < 100:



                continue



            for a1, b1 in Hy[ya]:



                for a2, b2 in Hy[yb]:



                    lo = max(a1, a2)



                    hi = min(b1, b2)



                    if hi - lo < 100:



                        continue



                    if v_span(lo, ya, yb) and v_span(hi, ya, yb):



                        rects.append((lo * SNAP, ya * SNAP,



                                      hi * SNAP, yb * SNAP))







    def area(r):



        return (r[2] - r[0]) * (r[3] - r[1])







    def contains(big, o):



        return (big[0] - 5 <= o[0] and big[1] - 5 <= o[1]



                and o[2] <= big[2] + 5 and o[3] <= big[3] + 5)







    # 仅保留极大矩形（去除被包含的）



    rects.sort(key=lambda r: -area(r))



    maximal = []



    for r in rects:



        if any(contains(f, r) for f in maximal):



            continue



        maximal.append(r)



    if not maximal:



        return 0







    # 若某个超大矩形包含 >=3 个其它矩形，则小矩形才是真正的图框



    containers = [r for r in maximal



                  if sum(1 for o in maximal if contains(r, o)) >= 3]



    if containers:



        big = max(containers, key=area)



        frames = [o for o in maximal if contains(big, o)]



    else:



        frames = maximal







    def sheet_like(r):



        w = r[2] - r[0]



        h = r[3] - r[1]



        if w < 150 or h < 150:



            return False



        if area(r) < 40000:



            return False



        ar = max(w, h) / min(w, h)



        return ar <= 10     # 排除极端细长的误判线条







    return len([r for r in frames if sheet_like(r)])











def _norm_sheet_code(val):



    """从图号/页码属性值中规整出『图纸代码』，用于去重统计张数。







    统一取 4~6 位连续数字；忽略前导字母(如 T/A)与后缀，



    使 'T03004' 与 '03004'、'T12001' 与 '12001' 归并为同一张图纸。



    """



    if not val:



        return None



    val = val.strip().upper()



    m = re.search(r"\d{4,6}", val)



    if m:



        return m.group(0)



    return None











def _distinct_sheet_numbers(doc):



    """从标题块属性提取不重复的图纸张数。







    规则：



    * 主标题块：带有图号/页码类属性，且带图名/标题属性、或块名像标题栏、



      或带强图号属性(图号/页号/SHEET_NUMBER/…) → 取其规整后的图号值计为一张图纸。



    * 局部标题块(大样/节点/详图…)：其 DWGNO/所在图 指向『所属图纸』，



      也并入图纸集合(去重后即该图纸张数)。



    * 匿名且仅带 DWGNO/NO 的块(多为详图/截面标记，NO 为细部编号)不计为独立图纸，



      避免把『详图』当『图纸』重复计数(旧版 181 页即此类误判)。



    * 幕墙/构件等(不带图号属性)天然被过滤(旧版因块名含 BORDER 误判



      'Left Border Mullion' 的问题不再出现)。



    返回去重后的图纸代码数量。



    """



    sheet_tag = re.compile(



        r"(SHEET_NUMBER|图号|页号|DWGNO|NO|PAGE|SHEET|DRAWING_NO|图纸编号|图纸|"



        r"SOUZAITU|SUOZAI|所在图)",



        re.I)



    strong_sheet = re.compile(



        r"(SHEET_NUMBER|图号|页号|DRAWING_NO|图纸编号)", re.I)



    title_tag = re.compile(r"(DRAWING_TITLE|图名|标题|TITLE|NAME)", re.I)



    title_name = re.compile(



        r"(图框|边框|幅面|图纸|标题栏|会签栏|签名栏|FRAME|FRM|BORDER|BORD|"



        r"TITLE|TBAR|TB|SHEET|FORMAT|GB|国标)", re.I)



    detail_name = re.compile(



        r"(大样|节点|DETAIL|CALL|CALLOUT|NOTE|局部|NOSING|标高|型材|SECTION|DET|"



        r"局部放大|索引|放大|详图|节点详图|大样图|详图索引)", re.I)



    sheets = set()



    for layout in doc.layouts:



        for e in layout.query("INSERT"):



            try:



                nm = e.dxf.name



            except Exception:



                continue



            is_detail = bool(detail_name.search(nm))



            vals = []



            has_title = False



            any_tag = []



            for attr in e.attribs:



                try:



                    tag = attr.dxf.tag



                    txt = (attr.dxf.text or "").strip()



                except Exception:



                    continue



                any_tag.append(tag)



                if title_tag.search(tag):



                    has_title = True



                if sheet_tag.search(tag) and txt:



                    vals.append(txt)



            if not vals:



                continue



            if is_detail:



                # 局部标题：其图号指向所属图纸



                for v in vals:



                    c = _norm_sheet_code(v)



                    if c:



                        sheets.add(c)



                continue



            # 主标题块：需带图名/标题属性，或块名像标题栏，或带强图号属性



            if has_title or title_name.search(nm) or strong_sheet.search(" ".join(any_tag)):



                for v in vals:



                    c = _norm_sheet_code(v)



                    if c:



                        sheets.add(c)



    return len(sheets)











def _raw_closed_polylines(dxf_path):



    """从原始 DXF 文本解析闭合 LWPOLYLINE/POLYLINE 矩形的外接框。







    用于 ezdxf 无法解析的损坏 DXF 的兜底页数估算。



    """



    try:



        with open(dxf_path, "rb") as f:



            raw = f.read()



    except Exception:



        return []



    try:



        lines = raw.decode("utf-8").split("\n")



    except Exception:



        lines = raw.decode("gb18030", "ignore").split("\n")



    n = len(lines)



    rects = []



    i = 0



    while i < n - 1:



        code = lines[i].strip()



        val = lines[i + 1].strip() if i + 1 < n else ""



        if code == "0" and val in ("LWPOLYLINE", "POLYLINE"):



            etype = val



            pts = []



            closed = False



            j = i + 2



            in_vertex = False



            while j < n - 1:



                c2 = lines[j].strip()



                v2 = lines[j + 1].strip() if j + 1 < n else ""



                if c2 == "0":



                    if etype == "POLYLINE":



                        if v2 == "VERTEX":



                            in_vertex = True



                            j += 2



                            continue



                        else:



                            break



                    else:



                        break



                if etype == "LWPOLYLINE":



                    if c2 == "70":



                        try:



                            if int(float(v2)) & 1:



                                closed = True



                        except Exception:



                            pass



                    elif c2 == "10":



                        try:



                            x = float(v2)



                            if j + 3 < n and lines[j + 2].strip() == "20":



                                y = float(lines[j + 3].strip())



                                pts.append((x, y))



                                j += 2



                        except Exception:



                            pass



                else:



                    if in_vertex:



                        if c2 == "70":



                            try:



                                if int(float(v2)) & 1:



                                    closed = True



                            except Exception:



                                pass



                        elif c2 == "10":



                            try:



                                x = float(v2)



                                if j + 3 < n and lines[j + 2].strip() == "20":



                                    y = float(lines[j + 3].strip())



                                    pts.append((x, y))



                                    j += 2



                            except Exception:



                                pass



                j += 1



            if closed and len(pts) >= 4:



                xs = [p[0] for p in pts]; ys = [p[1] for p in pts]



                rects.append((min(xs), min(ys), max(xs), max(ys)))



            i = j



        else:



            i += 1



    return rects











def _count_geom_frames(rects, min_side=150, min_area=40000, max_ar=10):



    """从矩形列表统计图纸张数：过滤过小/过扁，去掉单一整体外框，



    取互不重叠包含的最大矩形个数，并剔除明显偏小的标题栏/详图框。"""



    cand = []



    for (a, b, c, d) in rects:



        w = c - a; h = d - b



        if w < min_side or h < min_side:



            continue



        area = w * h



        if area < min_area:



            continue



        ar = max(w, h) / max(min(w, h), 1)



        if ar > max_ar:



            continue



        cand.append((a, b, c, d, area))



    if not cand:



        return 0



    cand.sort(key=lambda r: -r[4])



    # 去掉单一整体外框：若最大者明显大于次大者



    if len(cand) >= 2 and cand[0][4] > cand[1][4] * 4:



        cand = cand[1:]



    if not cand:



        return 0







    def contains(big, o):



        return (big[0] - 5 <= o[0] and big[1] - 5 <= o[1]



                and o[2] <= big[2] + 5 and o[3] <= big[3] + 5)







    # 取互不包含的最大矩形（过滤框内小框）



    cand.sort(key=lambda r: -(r[4]))



    maximal = []



    for r in cand:



        if any(contains(f, r) for f in maximal):



            continue



        maximal.append(r)



    if not maximal:



        return 0







    # 按面积自然聚类：最大断层以上的矩形视为真正图框，



    # 避免把标题栏/详图小框/装饰矩形等统计成图纸。



    maximal.sort(key=lambda r: -r[4])



    if len(maximal) >= 2:



        ratios = [maximal[i][4] / max(maximal[i+1][4], 1)



                  for i in range(len(maximal) - 1)]



        max_gap_idx = max(range(len(ratios)), key=lambda i: ratios[i])



        if ratios[max_gap_idx] >= 5:



            maximal = maximal[:max_gap_idx + 1]







    # 再按面积阈值去掉残余极小框



    max_area = max(r[4] for r in maximal)



    thr = max(min_area, max_area * 0.01)



    maximal = [r for r in maximal if r[4] >= thr]







    return len(maximal)











def _detect_lwpolyline_sheets(doc):



    """统计模型空间中类似图纸的闭合 LWPOLYLINE 矩形数量。"""



    rects = []



    for e in doc.modelspace().query("LWPOLYLINE"):



        try:



            if not e.closed:



                continue



            pts = list(e.get_points())



            if len(pts) < 4:



                continue



            xs = [p[0] for p in pts]; ys = [p[1] for p in pts]



            rects.append((min(xs), min(ys), max(xs), max(ys)))



        except Exception:



            pass



    return _count_geom_frames(rects)











def _count_detail_sheets(doc):



    """针对『大样/节点』类详图：按标题块位置聚类估算图纸页数（兜底）。"""



    from collections import defaultdict



    detail_pat = re.compile(



        r"(?:大样|节点|DETAIL|CALLOUT|详图|detail|节点详图|大样图|详图索引)", re.I)



    detail_pos = []



    for layout in doc.layouts:



        for e in layout.query("INSERT"):



            try:



                name = e.dxf.name



            except Exception:



                continue



            if not detail_pat.search(name):



                continue



            has_id = False



            for attr in e.attribs:



                try:



                    tag = attr.dxf.tag



                    if tag in ("DWGNO", "NO", "PAGE", "页号", "图号", "SHEET"):



                        has_id = True



                        break



                except Exception:



                    pass



            if not has_id:



                continue



            try:



                detail_pos.append((e.dxf.insert.x, e.dxf.insert.y))



            except Exception:



                pass



    n = len(detail_pos)



    if n < 2:



        return 0



    # 取模型空间最大的闭合 LWPOLYLINE 作为『图框条带』尺寸参考



    largest = None



    max_area = 0



    for e in doc.modelspace().query("LWPOLYLINE"):



        try:



            if not e.closed:



                continue



            pts = list(e.get_points())



            if len(pts) < 4:



                continue



            xs = [p[0] for p in pts]; ys = [p[1] for p in pts]



            w = max(xs) - min(xs); h = max(ys) - min(ys)



            area = w * h



            if area > max_area:



                max_area = area



                largest = (w, h)



        except Exception:



            pass



    if not largest:



        return 0



    w, h = largest



    if w < 1 or h < 1:



        return 0



    ar = max(w, h) / max(min(w, h), 1)



    if ar > 3:



        grid_size = min(w, h)



    else:



        grid_size = max(w, h)



    # 防止用超大整体外框做网格（>20000 视为异常，放弃聚类）



    if grid_size > 20000 or grid_size < 1:



        return 0



    cells = defaultdict(int)



    for x, y in detail_pos:



        cells[(round(x / grid_size), round(y / grid_size))] += 1



    return len(cells)











def _raw_layout_count(dxf_path):



    """从原始 DXF 文本中统计图纸空间布局数（兼容 ezdxf 无法解析时）。"""



    try:



        with open(dxf_path, "rb") as f:



            raw = f.read()



        text = None



        for enc in ("utf-8", "gbk", "gb18030"):



            try:



                text = raw.decode(enc)



                break



            except Exception:



                continue



        if text is None:



            text = raw.decode("utf-8", "ignore")



        matches = re.findall(r"\*Paper_Space(?:\d*)", text)



        unique = set()



        for m in matches:



            if m == "*Paper_Space":



                m = "*Paper_Space0"



            unique.add(m)



        return len(unique)



    except Exception:



        return 0











def count_cad_frames(dxf_path):



    """统计 CAD 图框数（页数）。返回 (frames:int|None, reason:str|None)。







    reason 标注本次页数的『统计口径』，供 GUI 状态列透明展示：



      * 布局计数        —— 含图元的图纸空间布局数（最权威，一张图=一个布局）



      * 标题块图号      —— 标题块属性中不重复图号/页码（主+局部标题块）



      * 详图聚类估算    —— 『大样/节点』详图按位置聚类（每张详图算一页，存近似）



      * 几何图框估算    —— LINE/LWPOLYLINE 重建图框矩形（无标题块/布局时兜底，存近似）



      * DXF损坏·几何估算 / DXF损坏·布局估算 / DXF损坏·按1页估算



                          —— ezdxf 无法解析损坏 DXF 时的原始文本几何/布局兜底







    判定优先级（取最代表『图纸张数』的信号）：



      1. 含图元的图纸空间布局数（不含 Model）；



      2. 标题块不重复图号 + 详图聚类（详图优先：det>=sheets 时取详图数）；



      3. 模型/图纸空间 LINE/LWPOLYLINE 重建出的图框矩形；



      4. 有图元则按 1 张计；



      5. ezdxf 无法解析时：原始 DXF 文本几何 / 布局数兜底。



    """



    try:



        import ezdxf



    except Exception as e:



        return None, "无法加载 ezdxf 统计图框: %s" % e







    doc = None



    try:



        doc = ezdxf.readfile(dxf_path)



    except Exception:



        try:



            san = sanitize_dxf(dxf_path)



            doc = ezdxf.readfile(san)



        except Exception:



            try:



                from ezdxf.recover import readfile



                doc, _ = readfile(san)



            except Exception:



                doc = None







    if doc:



        # 1) 图纸空间布局（含图元）



        paper = 0



        try:



            for layout in doc.layouts:



                if layout.name.upper() == "MODEL":



                    continue



                try:



                    if len(list(layout)) > 0:



                        paper += 1



                except Exception:



                    pass



        except Exception:



            paper = 0



        if paper >= 1:



            return paper, "布局计数"







        # 2) 标题块不重复图号 + 详图聚类



        #    口径（详图优先）：当『大样/节点』详图按位置聚类得到的张数



        #    >= 主/局部标题块的不重复图号时，视为『每张详图算一页』的详图集，



        #    以详图聚类数为页数（如幕墙/节点详图集 38 张）。



        try:



            sheets = _distinct_sheet_numbers(doc)



        except Exception:



            sheets = 0



        try:



            det = _count_detail_sheets(doc)



        except Exception:



            det = 0



        if det >= 1 and det >= sheets:



            return det, "详图聚类估算"



        if sheets >= 1:



            return sheets, "标题块图号"







        # 3) 几何图框矩形



        geo = 0



        try:



            geo_line = _detect_frame_rectangles(doc)



        except Exception:



            geo_line = 0



        try:



            geo_lw = _detect_lwpolyline_sheets(doc)



        except Exception:



            geo_lw = 0



        geo = max(geo_line, geo_lw)



        if geo >= 1:



            return geo, "几何图框估算"







        # 4) 大样/节点详图聚类（无图号时的兜底）



        if det >= 1:



            return det, "详图聚类估算"







        # 5) 有图元 -> 1



        try:



            if len(doc.modelspace()) > 0:



                return 1, "有图元·按1页估"



        except Exception:



            pass



        return None, "CAD 无图框/布局，无法统计页数"







    # 6) ezdxf 无法解析：原始 DXF 文本兜底



    rects = _raw_closed_polylines(dxf_path)



    geo = _count_geom_frames(rects)



    if geo >= 1:



        return geo, "DXF损坏·几何估算"



    rl = _raw_layout_count(dxf_path)



    if rl >= 1:



        return rl, "DXF损坏·布局估算"



    return 1, "DXF损坏·按1页估算"











def extract_cad(path, out_dir, converter, base):



    """Return (items, meta). items = list of paragraph strings (vector + OLE)."""



    os.makedirs(out_dir, exist_ok=True)



    if path.lower().endswith(".dwg"):



        raise RuntimeError("移动端不支持 .dwg（需 dwg2dxf 转换器，Android 无此二进制；请用 .dxf）")



    dxf_path = path







    text = extract_text_from_dxf(dxf_path, converter)



    items = [ln for ln in text.split("\n") if ln.strip()]







    emb = ""



    meta = {"ole_count": 0, "emb_objects": 0}



    try:



        import cad_ole_ocr as _olemod



        er = _olemod.extract_embedded_text(dxf_path, out_dir, base)



        emb = er.get("joined", "") or ""



        meta = {"ole_count": er.get("ole_count", 0),



                "emb_objects": er.get("unique_objects", 0)}



    except Exception:



        pass







    # 统计 CAD 图框数（标题块/布局），供 GUI『页数』列展示



    frames, freason = count_cad_frames(dxf_path)



    meta["pages"] = frames



    meta["pages_reason"] = freason







    # 检测：中文是否可能因 DWG→DXF 转换器(对中文 CJK 支持有限)而丢失。



    # 现象：源 dwg 二进制中含大量中文，但转换器提取到的中文极少。



    try:



        with open(path, "rb") as _f:



            _raw = _f.read()



        _src_cjk = sum(



            1 for c in _raw.decode("utf-16-le", "ignore")



            if "\u4e00" <= c <= "\u9fff"



        )



        _ext_cjk = 0



        for _it in items:



            _ext_cjk += sum(1 for c in _it if "\u4e00" <= c <= "\u9fff")



        if _src_cjk > 200 and _ext_cjk < max(10, _src_cjk * 0.05):



            meta["chinese_loss"] = True



            meta["chinese_loss_note"] = (



                "源文件含约 %d 个中文字符，但 DWG→DXF 转换器仅提取到 %d 个。"



                "中文很可能因转换器对中文(CJK)支持有限而丢失。"



                "建议：提供“文字型”PDF（CAD 另存为 PDF 时勾选“保留文字/作为文字而非曲线”），"



                "或用 TrueType 字体重新保存 DWG 后再统计，即可准确统计中文。"



                % (_src_cjk, _ext_cjk)



            )



    except Exception:



        pass







    for ln in emb.split("\n"):



        if ln.strip():



            items.append(ln)







    # 清理中间文件



    for junk in (dxf_path + "._sanitized.dxf",



                 os.path.splitext(dxf_path)[0] + "._min.dxf"):



        if os.path.exists(junk):



            try:



                os.remove(junk)



            except Exception:



                pass



    return items, meta











# ---------------------------------------------------------------------------



# 提取器：TXT



# ---------------------------------------------------------------------------



def extract_txt(path):



    with open(path, "r", encoding="utf-8", errors="replace") as f:



        text = f.read()



    items = []



    for ln in text.split("\n"):



        s = ln.strip()



        if s:



            items.append(s)



    return items











# ---------------------------------------------------------------------------



# 提取器：Word 文档 (.docx / .doc)



# ---------------------------------------------------------------------------



def extract_docx(path):



    from docx import Document



    doc = Document(path)



    items = []



    for p in doc.paragraphs:



        t = (p.text or "").strip()



        if t:



            items.append(t)



    for tbl in doc.tables:



        if tbl.rows:



            rows = [[(c.text if c.text is not None else "") for c in row.cells]



                    for row in tbl.rows]



            items.append(("table", rows))



    return items











def extract_doc(path):



    """旧版 .doc（二进制）用 Word COM 提取文字与表格。"""



    import pythoncom



    import win32com.client



    try:



        pythoncom.CoInitialize()



    except Exception:



        pass



    w = win32com.client.Dispatch("Word.Application")



    w.Visible = False



    w.DisplayAlerts = False



    try:



        w.AutomationSecurity = 1  # msoAutomationSecurityLow：跳过受保护视图



    except Exception:



        pass



    doc = w.Documents.Open(os.path.abspath(path))



    try:



        items = []



        for para in doc.Paragraphs:



            t = (para.Range.Text or "").replace("\x07", "").replace("\r", "").strip()



            if t:



                items.append(t)



        for ti in range(1, doc.Tables.Count + 1):



            tbl = doc.Tables(ti)



            if tbl.Rows.Count and tbl.Columns.Count:



                rows = [[(tbl.Cell(r, c).Range.Text or "")



                         .replace("\x07", "").replace("\r", "").strip()



                         for c in range(1, tbl.Columns.Count + 1)]



                        for r in range(1, tbl.Rows.Count + 1)]



                items.append(("table", rows))



        return items



    finally:



        doc.Close(False)



        w.Quit()











# ---------------------------------------------------------------------------



# 提取器：PowerPoint (.pptx / .ppt)



# ---------------------------------------------------------------------------



def extract_ppt(path):



    """旧版 .ppt（二进制）用 PowerPoint COM 提取文字与表格。"""



    import pythoncom



    import win32com.client



    try:



        pythoncom.CoInitialize()



    except Exception:



        pass



    p = win32com.client.Dispatch("PowerPoint.Application")



    try:



        pres = p.Presentations.Open(os.path.abspath(path), False, False, False)



        items = []



        for slide in pres.Slides:



            for shape in slide.Shapes:



                if shape.HasTextFrame:



                    t = (shape.TextFrame.TextRange.Text or "")



                    t = t.replace("\x0b", "").replace("\r", "").replace("\x07", "").strip()



                    if t:



                        items.append(t)



                if getattr(shape, "HasTable", False):



                    tbl = shape.Table



                    if tbl.Rows.Count and tbl.Columns.Count:



                        rows = [[(tbl.Cell(r, c).Text or "")



                                 .replace("\x0b", "").replace("\r", "").replace("\x07", "").strip()



                                 for c in range(1, tbl.Columns.Count + 1)]



                                for r in range(1, tbl.Rows.Count + 1)]



                        items.append(("table", rows))



        return items



    finally:



        pres.Close()



        p.Quit()











# ---------------------------------------------------------------------------



# 提取器：图片（OCR，RapidOCR onnxruntime，模型随包内置无需联网）



# ---------------------------------------------------------------------------



_OCR_ENGINE = None







def _get_ocr():



    global _OCR_ENGINE



    if _OCR_ENGINE is None:



        try:



            from rapidocr_onnxruntime import RapidOCR



        except Exception as e:



            raise RuntimeError(



                "OCR 组件未安装（缺失 opencv-python / onnxruntime / rapidocr_onnxruntime）。"



                "请确认构建时已在 chaquopy 的 pip 中声明这些依赖。"



            )



        d = os.environ.get("WORDCOUNT_OCR_DIR")



        if d and os.path.isdir(d):



            # 移动端：OCR 模型由 App 从 GitHub Release 下载后放在此目录



            det = os.environ.get("WORDCOUNT_OCR_DET") or os.path.join(d, "det", "ch_PP-OCRv4_det_infer.onnx")



            rec = os.environ.get("WORDCOUNT_OCR_REC") or os.path.join(d, "rec", "ch_PP-OCRv4_rec_infer.onnx")



            cls = os.environ.get("WORDCOUNT_OCR_CLS") or os.path.join(d, "cls", "ch_ppocr_mobile_v2.0_cls_infer.onnx")



            _OCR_ENGINE = RapidOCR(det_model_path=det, rec_model_path=rec, cls_model_path=cls)



        else:



            _OCR_ENGINE = RapidOCR()



    return _OCR_ENGINE











def extract_image(path):



    ocr = _get_ocr()



    result, _ = ocr(path)



    items = []



    if result:



        for box, text, score in result:



            t = (text or "").strip()



            if t:



                items.append(t)



    return items











# ---------------------------------------------------------------------------



# 压缩文件：自动解压 -> 递归收集内部受支持文件（支持嵌套压缩包）



# ---------------------------------------------------------------------------



# 注意：更长的后缀（.tar.gz）要排在前面，避免被 .gz 误判为单文件压缩



ARCHIVE_EXTS = (".zip", ".rar", ".7z", ".tar",



                ".tar.gz", ".tgz", ".tar.bz2", ".tbz2", ".tar.xz", ".txz",



                ".gz", ".bz2", ".xz")











def _archive_kind(path):



    low = path.lower()



    if low.endswith(".zip"):



        return "zip"



    if low.endswith(".rar"):



        return "rar"



    if low.endswith(".7z"):



        return "7z"



    for e in (".tar.gz", ".tgz", ".tar.bz2", ".tbz2", ".tar.xz", ".txz", ".tar"):



        if low.endswith(e):



            return "tar"



    for e in (".gz", ".bz2", ".xz"):



        if low.endswith(e):



            return "single"   # 单文件压缩，如 report.pdf.gz



    return None











def is_archive(path):



    return _archive_kind(path) is not None











def _safe_join(dest, name):



    """防目录穿越：确保解出的路径仍在 dest 之内。"""



    target = os.path.normpath(os.path.join(dest, name))



    if not os.path.abspath(target).startswith(os.path.abspath(dest)):



        return None



    return target











def _extract_zip(path, dest):



    with zipfile.ZipFile(path) as z:



        for info in z.infolist():



            name = info.filename



            # ZIP 若未置 UTF-8 标志(0x800)，文件名按 cp437 存储，中文需转 GBK/UTF-8



            if not (info.flag_bits & 0x800):



                for enc in ("gbk", "utf-8"):



                    try:



                        name = info.filename.encode("cp437").decode(enc)



                        break



                    except Exception:



                        continue



            target = _safe_join(dest, name)



            if target is None:



                continue



            if name.endswith("/") or info.is_dir():



                os.makedirs(target, exist_ok=True)



                continue



            os.makedirs(os.path.dirname(target), exist_ok=True)



            with z.open(info) as src, open(target, "wb") as dst:



                shutil.copyfileobj(src, dst)











def _extract_single(path, dest):



    """.gz/.bz2/.xz 单文件压缩：解出去掉压缩后缀的原文件。"""



    import gzip



    import bz2



    import lzma



    low = path.lower()



    base = os.path.basename(path)



    if low.endswith(".gz"):



        opener, out = gzip.open, base[:-3]



    elif low.endswith(".bz2"):



        opener, out = bz2.open, base[:-4]



    else:  # .xz



        opener, out = lzma.open, base[:-3]



    if not out:



        out = base + ".out"



    out_path = os.path.join(dest, out)



    with opener(path, "rb") as src, open(out_path, "wb") as dst:



        shutil.copyfileobj(src, dst)











def find_7z():



    """查找 7z.exe 解压工具（优先打包内置，再工具目录，再 PATH）。"""



    candidates = []



    if hasattr(sys, "_MEIPASS"):



        candidates.append(os.path.join(sys._MEIPASS, "7z", "7z.exe"))



    candidates.append(os.path.join(os.path.dirname(os.path.abspath(__file__)), "7z", "7z.exe"))



    candidates.append(os.path.join(_HOME, ".workbuddy", "tools", "7z", "7z.exe"))



    for base in os.environ.get("PATH", "").split(os.pathsep):



        if base:



            candidates.append(os.path.join(base, "7z.exe"))



    for p in candidates:



        if os.path.exists(p):



            return p



    return None











_SEVENZ = None



def _get_7z():



    global _SEVENZ



    if _SEVENZ is None:



        _SEVENZ = find_7z()



    return _SEVENZ











def _extract_with_7z(path, dest):



    """调用 7z.exe 解压。path 和 dest 可含空格。"""



    tool = _get_7z()



    if not tool:



        raise RuntimeError("未找到 7z.exe 解压工具，无法处理 .rar/.7z/.tar 等格式")



    # 7z 的 -o 参数不能有空格，-aoa 覆盖模式，-y 全部确认，-bd 禁止进度条



    cmd = [tool, "x", path, "-o" + dest, "-y", "-bd", "-aoa"]



    r = subprocess.run(cmd, capture_output=True, timeout=300)



    if r.returncode != 0:



        err = (r.stderr.decode("gbk", "ignore") or r.stdout.decode("gbk", "ignore") or "").strip()



        raise RuntimeError("7z 解压失败 (rc=%d): %s" % (r.returncode, err[:200]))











def extract_archive(path, dest):



    """把压缩包解到 dest 目录。返回 dest。







    策略：



    - .zip 优先用 Python 标准库处理，以正确恢复中文文件名编码；



    - 其他格式（.rar/.7z/.tar/.gz/.bz2/.xz）优先用 7z.exe 外部工具；



    - 未找到 7z.exe 时回退到 Python 库（.rar 回退需要本机 unrar/bsdtar）。



    """



    os.makedirs(dest, exist_ok=True)



    kind = _archive_kind(path)



    if kind == "zip":



        _extract_zip(path, dest)



    elif kind == "tar":



        if _get_7z():



            _extract_with_7z(path, dest)



        else:



            with tarfile.open(path, "r:*") as tf:



                try:



                    tf.extractall(dest, filter="data")   # Py3.12+ 安全过滤



                except TypeError:



                    tf.extractall(dest)



    elif kind == "rar":



        if _get_7z():



            _extract_with_7z(path, dest)



        else:



            import rarfile



            with rarfile.RarFile(path) as rf:



                rf.extractall(dest)



    elif kind == "7z":



        if _get_7z():



            _extract_with_7z(path, dest)



        else:



            import py7zr



            with py7zr.SevenZipFile(path, "r") as z:



                z.extractall(dest)



    elif kind == "single":



        if _get_7z():



            _extract_with_7z(path, dest)



        else:



            _extract_single(path, dest)



    else:



        raise RuntimeError("不支持的压缩格式: " + os.path.basename(path))



    return dest











def _gather_supported_from_dir(root, base_root, _depth=0):



    """遍历目录取所有受支持文件；遇到嵌套压缩包则继续解压（最多 5 层防炸弹）。



    返回 [(abspath, relname), ...]，relname 为相对 base_root（顶层解压目录）的相对路径，



    用于作为压缩包内每个文件的展示名（保留目录层级，如 outer/nested.zip/file.docx）。"""



    out = []



    for dirpath, _, files in os.walk(root):



        for fn in files:



            fp = os.path.join(dirpath, fn)



            if fn.lower().endswith(SUPPORTED):



                rel = os.path.relpath(fp, base_root)



                out.append((fp, rel))



            elif is_archive(fp) and _depth < 5:



                sub = fp + "._wc_extracted"



                try:



                    extract_archive(fp, sub)



                    out += _gather_supported_from_dir(sub, base_root, _depth + 1)



                except Exception as e:



                    print(f"[WARN] 嵌套解压失败 {fn}: {e}")



    return out











def expand_archive_file(path):



    """解压一个压缩包，返回 [(abspath, relname), ...]（含嵌套压缩包内部文件）。"""



    dest = tempfile.mkdtemp(prefix="wc_arc_")



    extract_archive(path, dest)



    return _gather_supported_from_dir(dest, dest)











def count_archive(path):



    """解压压缩包，统计内部所有受支持文件，返回 (combined_items, summed_stats, meta)。







    combined_items：内部所有文件提取出的文字合并（用于合并 Word 报告）。



    summed_stats：汇总统计（压缩包这一行/条目显示的总数）。



    meta：



        inner_count：成功统计的内部文件数（不含统计失败的）。



        inner：每个内部文件的明细列表，供 GUI 展开查看（每个含独立 stats/meta）。



        pages：内部可统计页数文件的页数合计，供汇总行累加。



        _extract_dir：解压目录（保留不删），供后续导出渲染内部文件时直接访问。



    """



    dest = tempfile.mkdtemp(prefix="wc_arc_")



    extract_archive(path, dest)



    inner = _gather_supported_from_dir(dest, dest)



    combined = []



    stats = {"words": 0, "fe": 0, "nc": 0, "chars": 0}



    cnt = 0



    total_pages = 0



    inner_details = []



    for f, rel in inner:



        try:



            tmp = tempfile.mkdtemp(prefix="wc_ai_")



            items, _meta = extract_for(f, "all", tmp, None, False)



            combined += items



            s = count_items(items)



            for k in stats:



                stats[k] += s[k]



            pg = _meta.get("pages")



            if isinstance(pg, int):



                total_pages += pg



            inner_details.append({"name": rel, "stats": s,



                                  "meta": _meta})



            cnt += 1



        except Exception:



            pass



    meta = {"inner_count": cnt, "inner": inner_details, "_extract_dir": dest}



    if cnt:



        # 压缩包行页数列显示内部文件页数合计；无合计时显示说明



        if total_pages > 0:



            meta["pages"] = total_pages



            meta["pages_reason"] = None



        else:



            meta["pages"] = None



            meta["pages_reason"] = "压缩包(内部 %d 个文件)" % cnt



    return combined, stats, meta











# ---------------------------------------------------------------------------



# 格式分发



# ---------------------------------------------------------------------------



SUPPORTED = (".pdf", ".dwg", ".dxf", ".xlsx", ".xlsm", ".pptx", ".ppt",



             ".docx", ".doc", ".txt", ".ai", ".cdr",



             ".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".gif", ".webp")











def count_pages(ext, path):



    """统计文件页数。返回 (pages:int|None, reason:str|None)。







    能统计的给出整数页数；无法统计（如表格文件无页码概念）给出原因，



    供 GUI『页数』列在统计失败时显示原因。



    """



    try:



        if ext == ".pdf":



            import fitz



            d = fitz.open(path)



            n = len(d)



            d.close()



            return n, None



        if ext in (".xlsx", ".xlsm"):



            import openpyxl



            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)



            n = 0



            for ws in wb:



                has_text = False



                for row in ws.iter_rows(values_only=True):



                    for v in row:



                        if v is not None and str(v).strip():



                            n += 1



                            has_text = True



                            break



                    if has_text:



                        break



            if n:



                return n, None



            return 0, "所有工作表均无文字"



        if ext == ".pptx":



            from pptx import Presentation



            return len(Presentation(path).slides), None



        if ext == ".ppt":



            import pythoncom



            import win32com.client



            try:



                pythoncom.CoInitialize()



            except Exception:



                pass



            p = win32com.client.Dispatch("PowerPoint.Application")



            try:



                pres = p.Presentations.Open(os.path.abspath(path), False, False, False)



                return pres.Slides.Count, None



            finally:



                try:



                    pres.Close()



                except Exception:



                    pass



                try:



                    p.Quit()



                except Exception:



                    pass



        if ext == ".doc":



            import pythoncom



            import win32com.client



            try:



                pythoncom.CoInitialize()



            except Exception:



                pass



            w = win32com.client.Dispatch("Word.Application")



            w.Visible = False



            try:



                doc = w.Documents.Open(os.path.abspath(path))



                return doc.ComputeStatistics(2), None  # wdStatisticPages = 2



            finally:



                try:



                    doc.Close(False)



                except Exception:



                    pass



                try:



                    w.Quit()



                except Exception:



                    pass



        if ext == ".docx":



            return None, ".docx 无页码信息(需 Word 才能统计)"



        if ext == ".txt":



            return 1, None



        if ext in (".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".gif", ".webp"):



            return 1, None



        if ext in (".dwg", ".dxf"):



            return None, None  # CAD 图框由 extract_cad 单独统计



        return None, "不支持的格式，无法统计页数"



    except Exception as e:



        return None, "页数统计失败: %s" % e











def extract_for(path, sheet_filter, out_dir, converter, with_notes):



    ext = os.path.splitext(path)[1].lower()



    base = os.path.splitext(os.path.basename(path))[0]



    if ext == ".pdf":



        items, meta = extract_pdf(path)



        pg, reason = count_pages(ext, path)



        meta["pages"] = pg



        meta["pages_reason"] = reason



        return items, meta



    if ext in (".xlsx", ".xlsm"):



        items, sheets = extract_excel(path, sheet_filter)



        pg, reason = count_pages(ext, path)



        return items, {"pages": pg, "pages_reason": reason, "sheets": sheets}



    if ext == ".pptx":



        items = extract_pptx(path, with_notes)



        pg, reason = count_pages(ext, path)



        return items, {"pages": pg, "pages_reason": reason}



    if ext == ".ppt":



        raise RuntimeError("移动端不支持 .ppt（需 Windows PowerPoint COM；请用 .pptx）")



    if ext in (".dwg", ".dxf"):



        # extract_cad 已在 meta 中写入 pages(图框)/pages_reason



        return extract_cad(path, out_dir, converter, base)



    if ext == ".docx":



        items = extract_docx(path)



        pg, reason = count_pages(ext, path)



        return items, {"pages": pg, "pages_reason": reason}



    if ext == ".doc":



        raise RuntimeError("移动端不支持 .doc（需 Windows Word COM；请用 .docx）")



    if ext == ".txt":



        items = extract_txt(path)



        pg, reason = count_pages(ext, path)



        return items, {"pages": pg, "pages_reason": reason}



    if ext in (".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".gif", ".webp"):



        items = extract_image(path)



        pg, reason = count_pages(ext, path)



        return items, {"pages": pg, "pages_reason": reason}



    if ext == ".ai":



        # 设计图：单页，meta 内已置 pages=1



        return extract_ai(path)



    if ext == ".cdr":



        # 设计图：单页，meta 内已置 pages=1



        return extract_cdr(path)



    raise RuntimeError("不支持的格式: %s" % ext)











# ---------------------------------------------------------------------------



# 输入收集



# ---------------------------------------------------------------------------



def collect_inputs(inputs, recursive):



    files = []



    archives = []



    for it in inputs:



        if os.path.isdir(it):



            pats = ["**/*" + e for e in SUPPORTED] if recursive else ["*" + e for e in SUPPORTED]



            for pat in pats:



                files += glob.glob(os.path.join(it, pat), recursive=recursive)



            apats = ["**/*" + e for e in ARCHIVE_EXTS] if recursive else ["*" + e for e in ARCHIVE_EXTS]



            for pat in apats:



                archives += glob.glob(os.path.join(it, pat), recursive=recursive)



        elif os.path.isfile(it):



            if is_archive(it):



                archives.append(it)



            else:



                files.append(it)



        else:



            for g in glob.glob(it):



                if is_archive(g):



                    archives.append(g)



                elif os.path.isfile(g):



                    files.append(g)



    # 展开压缩包（按绝对路径去重，避免 .tar.gz 同时匹配 *.gz 造成重复解压）



    for arc in dict.fromkeys(os.path.abspath(a) for a in archives):



        try:



            files += [f for f, _ in expand_archive_file(arc)]



        except Exception as e:



            print(f"[ERR] 解压失败 {os.path.basename(arc)}: {e}")



    seen, out = set(), []



    for f in files:



        fl = f.lower()



        if fl.endswith(SUPPORTED) and f not in seen:



            seen.add(f)



            out.append(f)



    return sorted(out)











# ---------------------------------------------------------------------------



# 发送（qqmail-send-file 两阶段）



# ---------------------------------------------------------------------------



def prepare_send(out_dir, results, png_path, to, subject, body, merged_path=None):



    zip_path = os.path.join(out_dir, "wordcount_batch.zip")



    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as z:



        for r in results:



            if r.get("docx") and os.path.exists(r["docx"]):



                z.write(r["docx"], os.path.basename(r["docx"]))



        if merged_path and os.path.exists(merged_path):



            z.write(merged_path, os.path.basename(merged_path))



        if png_path and os.path.exists(png_path):



            z.write(png_path, os.path.basename(png_path))



    size = os.path.getsize(zip_path)



    print(f"\n[打包] {zip_path} ({size} 字节)")



    if size > 1_000_000:



        print("[WARN] 超过 QQ 邮箱单附件 1MB 上限，跳过发送。请手动分发产物。")



        return



    send_script = os.path.expanduser(



        "~/.workbuddy/skills/qqmail-send-file/scripts/send_file.py")



    if not os.path.exists(send_script):



        print("[WARN] 未找到 qqmail-send-file 技能，无法发送。")



        return



    subj = subject or "字数统计（Word口径）批量结果"



    body = body or "附件为批量字数统计结果（口径同 Word 字数统计对话框）。"



    cmd = [sys.executable, send_script, "--file", zip_path,



           "--to", to, "--subject", subj, "--body", body]



    try:



        out = subprocess.run(cmd, capture_output=True, text=True, timeout=120)



        sys.stdout.write(out.stdout)



        if out.stderr:



            sys.stderr.write(out.stderr)



        try:



            s = out.stdout[out.stdout.find("{"):out.stdout.rfind("}") + 1]



            j = json.loads(s)



            tok = j.get("confirmation_token")



            if tok:



                print(f"\n[发送准备完成] 确认令牌: {tok}")



                print("请运行 qqmail-send-file 的确认步骤（--confirm-token <令牌>）完成发送。")



        except Exception:



            pass



    except Exception as e:



        print(f"[ERR] 发送准备失败: {e}")











# ---------------------------------------------------------------------------



# 主流程



# ---------------------------------------------------------------------------



def _unique_path(path):



    """返回不覆盖已有文件的目标路径（同名时追加 _1/_2 ...）。"""



    if not os.path.exists(path):



        return path



    root, ext = os.path.splitext(path)



    i = 1



    while os.path.exists(f"{root}_{i}{ext}"):



        i += 1



    return f"{root}_{i}{ext}"











# ---------------------------------------------------------------------------



# 导出「无法准确统计的内容」



#   CAD 图纸中的文字、图片型 PDF / PDF 图片页中的文字，统计口径可能不准。



#   把这些内容渲染成图片（CAD 按布局截图、PDF 提取整页图片），融合为一个



#   PDF 导出，每页标注：文件名 / 第几页(布局或页码) / 本程序统计的字数·中文·非中文。



# ---------------------------------------------------------------------------



def _safe_name(s):



    return re.sub(r'[\\/:*?"<>|\r\n\t]', "_", str(s))[:50]











def _cad_layout_pages(dxf_path, work_dir, converter):



    """渲染 CAD(dwg/dxf) 的每个布局（模型空间 + 图纸空间）为 PNG 页。







    采用 matplotlib 后端：对坏坐标实体、空块名 INSERT、二进制数据损坏等



    容错性远优于 pymupdf 后端（后者在部分 DWG 上会硬崩溃，无法被 Python 捕获）。



    返回 [(png_path, label), ...]；若无法渲染返回 []。"""



    import os as _os



    try:



        import ezdxf



        import matplotlib



        matplotlib.use("Agg")



        import matplotlib.pyplot as plt



        from ezdxf.addons.drawing import frontend, layout as dlayout, matplotlib as ezmpl



    except Exception:



        # 移动端未内置 matplotlib / CAD 截图渲染，跳过 CAD 截图（图片型 PDF 仍可导出）



        return []







    # DWG 需先经 dwg2dxf 转为 DXF（ezdxf 不能直接读 DWG）



    render_src = dxf_path



    if dxf_path.lower().endswith(".dwg"):



        conv = find_converter(converter)



        if not conv:



            return []



        dxf_out = _os.path.join(



            work_dir,



            _safe_name(_os.path.splitext(_os.path.basename(dxf_path))[0]) + ".dxf")



        try:



            dwg_to_dxf(dxf_path, dxf_out, conv)



        except Exception:



            return []



        render_src = dxf_out







    # 容错读取 DXF：标准读 → 净化(去二进制坏数据) → recover



    doc = None



    try:



        doc = ezdxf.readfile(render_src)



    except Exception:



        try:



            san = sanitize_dxf(render_src)



            doc = ezdxf.readfile(san)



        except Exception:



            try:



                from ezdxf.recover import readfile as _recover



                doc, _ = _recover(san)



            except Exception:



                doc = None



    if not doc:



        return []







    def filt(e):



        # 跳过块名为空的 INSERT（引用了不存在的块定义，会导致展开虚拟实体崩溃）



        if e.dxftype() == "INSERT":



            nm = (getattr(e.dxf, "name", "") or "").strip()



            if not nm:



                return False



            try:



                if doc.blocks.get(nm) is None:



                    return False



            except Exception:



                return False



        return True







    pages = []



    prefix = _safe_name(_os.path.splitext(_os.path.basename(render_src))[0])



    for lo in doc.layouts:



        name = lo.name



        try:



            fig = plt.figure(figsize=(16, 11))



            ax = fig.add_axes([0, 0, 1, 1])



            ctx = frontend.RenderContext(doc)



            backend = ezmpl.MatplotlibBackend(ax)



            fe = frontend.Frontend(ctx, backend)



            fe.draw_layout(lo, filter_func=filt)



            # 文件名带源文件前缀，避免多个 CAD 渲染到同一 work_dir 时互相覆盖



            png = _os.path.join(work_dir, "cad_%s_%s.png" % (prefix, _safe_name(name)))



            fig.savefig(png, dpi=110)



            plt.close(fig)



            label = "模型空间" if name.upper() == "MODEL" else ("图纸空间: " + name)



            pages.append((png, label))



        except Exception:



            continue



    # 释放 matplotlib / 大几何内存，避免连续渲染多个大 CAD 时内存累积导致崩溃



    try:



        plt.close("all")



    except Exception:



        pass



    try:



        import gc



        gc.collect()



    except Exception:



        pass



    return pages











def _pdf_unreliable_pages(pdf_path, work_dir, meta):



    """对图片型 PDF / 含图片页的 PDF，把『无法准确统计』的页渲染为 PNG。



    返回 [(png_path, label), ...]。"""



    import fitz



    img_pages = meta.get("img_pages") or []



    image_only = meta.get("image_only")



    try:



        d = fitz.open(pdf_path)



    except Exception:



        return []



    n = len(d)



    if image_only:



        targets = list(range(1, n + 1))



        tag = "整页图片·OCR"



    elif img_pages:



        targets = sorted(set(int(p) for p in img_pages))



        tag = "图片页·OCR"



    else:



        targets = []



    pages = []



    for pno in targets:



        try:



            sp = d[pno - 1]



            png = os.path.join(



                work_dir,



                "pdf_%s_%d.png"



                % (_safe_name(os.path.splitext(os.path.basename(pdf_path))[0]), pno))



            pix = sp.get_pixmap(matrix=fitz.Matrix(2, 2))



            pix.save(png)



            pages.append((png, "第 %d 页(%s)" % (pno, tag)))



        except Exception:



            continue



    d.close()



    return pages











def prepare_unreliable_entries(files_info, work_dir, converter=None):



    """准备「无法准确统计的内容」导出数据。







    files_info: list of (name, stats_dict, meta_dict, src_path, ext)



        name      : 展示用文件名（可含 压缩包/内部路径）



        stats_dict: 该文件统计 {words,fe,nc,chars}



        meta_dict : 该文件 meta（含 image_only / img_pages / pages 等）



        src_path  : 实际文件路径（CAD 需可访问；压缩包内文件为解压后的路径）



        ext       : 小写扩展名



    返回 entries: list of (name, stats_dict, pages)



        pages: list of (png_path, page_label)



    仅包含确有不可靠内容的文件。"""

    import os as _os

    # Chaquopy 传入的 files_info 可能是 Java ArrayList，递归转为原生 list
    files_info = _to_py_list(files_info)

    entries = []


    for (name, stats, meta, src, ext) in files_info:



        if not src or not _os.path.exists(src):



            continue



        if ext in (".dwg", ".dxf"):



            pages = _cad_layout_pages(src, work_dir, converter)



        elif ext == ".pdf":



            if meta.get("image_only") or meta.get("img_pages"):



                pages = _pdf_unreliable_pages(src, work_dir, meta)



            else:



                pages = []



        else:



            pages = []



        if pages:



            entries.append((name, stats, pages))



        # 每处理完一个文件释放内存（CAD 渲染体量大，连续渲染易内存累积）



        try:



            import gc



            gc.collect()



        except Exception:



            pass



    return entries











def build_unreliable_pdf(entries, out_path):



    """把 entries 合并为一个 PDF：每页 = 页头标注 + 内容图。



    页头：文件名 / 页标签 / 本程序统计 → 字数·中文·非中文。



    内容图统一缩放到固定宽度，页头字号随宽度缩放，保证标注始终可读。"""



    import fitz



    doc = fitz.open()



    total = 0



    TW = 1240  # 页面内容宽度（点），CAD 截图/PDF 页均缩放到此宽度



    for (name, stats, pages) in entries:



        words = int(stats.get("words", 0) or 0)



        fe = int(stats.get("fe", 0) or 0)



        nc = int(stats.get("nc", 0) or 0)



        for (png, label) in pages:



            try:



                ip = fitz.open(png)



                ipix = ip[0]



                iw, ih = ipix.rect.width, ipix.rect.height



                ip.close()



            except Exception:



                continue



            if iw <= 0 or ih <= 0:



                continue



            scale = TW / iw



            draw_w = TW



            draw_h = ih * scale



            font = max(12, int(TW / 90))



            head_h = font + 16



            page = doc.new_page(width=draw_w, height=head_h + draw_h)



            header = ("文件: %s    %s    本程序统计 → 字数 %d | 中文 %d | 非中文 %d"



                      % (name, label, words, fe, nc))



            page.insert_text((10, font + 4), header, fontsize=font, color=(0, 0, 0))



            page.insert_image(fitz.Rect(0, head_h, draw_w, head_h + draw_h),



                              filename=png)



            total += 1



    doc.save(out_path)



    doc.close()



    return total











# ---------------------------------------------------------------------------



# 移动端 API（供 Chaquopy / Kotlin 调用，返回 JSON 可序列化 dict）



# ---------------------------------------------------------------------------



def count_file(path, sheet_filter="all", with_notes=False, converter=None):



    """统计单个文件，返回 {name, ext, is_archive, stats, meta, ...}。







    不支持的格式（.doc/.ppt/.dwg）直接抛 RuntimeError，由 Kotlin 层捕获并提示。



    """



    ext = os.path.splitext(path)[1].lower()



    if ext in (".doc", ".ppt"):



        raise RuntimeError("移动端不支持 %s（需 Windows Word/PowerPoint）" % ext)



    if ext == ".dwg":



        raise RuntimeError("移动端不支持 .dwg（需 dwg2dxf；请用 .dxf）")



    tmp = tempfile.mkdtemp(prefix="wc_")



    if is_archive(path):



        combined, stats, meta = count_archive(path)



        return {



            "name": os.path.basename(path), "ext": ext, "is_archive": True,



            "stats": stats, "meta": meta, "inner": meta.get("inner", []),



        }



    items, meta = extract_for(path, sheet_filter, tmp, converter, with_notes)



    stats = count_items(items)



    meta["ext"] = ext



    return {



        "name": os.path.basename(path), "ext": ext, "is_archive": False,



        "stats": stats, "meta": meta,



        "sheets": meta.get("sheets"),



        "pages": meta.get("pages"),



        "pages_reason": meta.get("pages_reason"),



        "img_pages": meta.get("img_pages"),



        "image_only": meta.get("image_only"),



    }











def _to_py_list(obj):
    """把 Chaquopy 从 Kotlin 传入的 Java 集合转成真正的 Python list。

    Chaquopy 在此项目中不会自动把 Kotlin List 转成 Python list，
    而是把原始 Java ArrayList 代理对象直接传进来。对它调用 Python 的
    list() 会触发 "ArrayList object is not iterable"，必须用 size()/get(i) 手动取出。
    嵌套列表（如 build_export_pdf 的 filesInfo）递归处理；字符串/数字/Map 等原子值原样返回。
    """
    # 已经是 Python list/tuple
    if isinstance(obj, (list, tuple)):
        return [_to_py_list(x) for x in obj]
    # Java List/ArrayList：用 size()/get(i) 取出（排除 Map：Map 也有 size/get 但带 put）
    if (not isinstance(obj, (str, bytes, dict))
            and hasattr(obj, "size") and hasattr(obj, "get")
            and not hasattr(obj, "put")):
        try:
            n = obj.size()
            return [_to_py_list(obj.get(i)) for i in range(n)]
        except Exception:
            pass
    # 原子值（字符串/数字/Map/dict 等）：原样返回，不要当可迭代对象拆开
    return obj


def count_files(paths, sheet_filter="all", with_notes=False):



    """批量统计多个文件，逐个容错，返回 [{ok, result|error, name}, ...]。"""


    # Chaquopy 从 Kotlin 传入的 List<String> 是 Java ArrayList 代理对象，
    # 不能用 list() 直接迭代，必须用 size()/get(i) 取出。
    paths = _to_py_list(paths)

    out = []



    for p in paths:



        try:



            out.append({"ok": True, "result": count_file(p, sheet_filter, with_notes),



                        "name": os.path.basename(p)})



        except Exception as e:
            out.append({"ok": False, "error": "%s\n%s" % (e, traceback.format_exc()),
                        "name": os.path.basename(p)})



    return out











def count_text(text, name="图片"):



    """移动端图片 OCR 在 Kotlin 层用 Tesseract 完成，识别出的文字传回此函数计数。







    返回结构同 count_file，便于 Kotlin 层统一解析（含 stats/meta）。



    """



    items = [ln.strip() for ln in (text or "").split("\n") if ln.strip()]



    stats = count_items(items)



    ext = os.path.splitext(name)[1].lower() or ".image"

    is_img = ext in (".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".gif", ".webp")

    meta = {"ext": ext, "image_ocr": is_img}



    return {



        "name": name, "ext": ext, "is_archive": False,



        "stats": stats, "meta": meta,



        "sheets": None, "pages": None, "pages_reason": ("图片(已OCR识别)" if is_img else None),



        "img_pages": None, "image_only": False,



    }











def build_export_pdf(files_info, out_path, work_dir=None, converter=None):
    # Chaquopy 传入的 filesInfo 是 Java List<List<Any?>>，递归转成 Python list
    files_info = _to_py_list(files_info)



    """导出『无法准确统计内容』PDF。







    files_info: list of (name, stats_dict, meta_dict, src_path, ext)



    返回生成的 PDF 路径；无可导出内容返回 None。



    """



    wd = work_dir or tempfile.mkdtemp(prefix="wc_exp_")



    entries = prepare_unreliable_entries(files_info, wd, converter)



    if not entries:



        return None



    build_unreliable_pdf(entries, out_path)



    return out_path











# ---------------------------------------------------------------------------



# 命令行入口（桌面端保留；移动端不调用）



# ---------------------------------------------------------------------------



def main():



    reexec_with_venv()



    ap = argparse.ArgumentParser(



        description="统一字数统计 v1.3.15（PDF/CAD/Excel/PPT -> Word 口径）")



    ap.add_argument("--version", action="version", version="%(prog)s v1.3.15")

    ap.add_argument("inputs", nargs="+", help="文件或目录（可多个）")



    ap.add_argument("--output", default=None, help="输出目录")



    ap.add_argument("--sheet", default="all", help="Excel: all | first | 工作表名")



    ap.add_argument("--recursive", action="store_true", help="目录递归")



    ap.add_argument("--with-notes", action="store_true", help="PPT 含备注文字")



    ap.add_argument("--summary-png", action="store_true", help="生成汇总 PNG")



    ap.add_argument("--merge", action="store_true",



                    help="合并为一个 Word 文档（含各文件内容 + 汇总表 + 汇总图，图在末尾）")



    ap.add_argument("--no-docx", action="store_true", help="不生成 Word，仅统计")



    ap.add_argument("--send", action="store_true", help="打包并发 QQ 邮箱")



    ap.add_argument("--to", default="oliverzh2020@qq.com")



    ap.add_argument("--subject", default=None)



    ap.add_argument("--body", default=None)



    ap.add_argument("--converter", default=None, help="dwg2dxf.exe 路径")



    args = ap.parse_args()







    out_dir = args.output or os.path.join(os.getcwd(), "wordcount_output")



    os.makedirs(out_dir, exist_ok=True)







    files = collect_inputs(args.inputs, args.recursive)



    if not files:



        print("未找到任何受支持的文件：", SUPPORTED)



        sys.exit(1)







    results = []



    for f in files:



        try:



            items, meta = extract_for(f, args.sheet, out_dir, args.converter, args.with_notes)



            base = os.path.splitext(os.path.basename(f))[0]



            docx_path = None



            stats = count_items(items)



            if not args.no_docx and not args.merge:



                docx_path = _unique_path(os.path.join(out_dir, base + "_字数统计.docx"))



                build_docx(base, items, docx_path, stats, meta)



            results.append({"src": f, "items": items,



                            "docx": docx_path if not args.no_docx else None,



                            "stats": stats, "meta": meta})



            s = stats



            img_tag = ""



            if meta.get("img_pages"):



                n = len(meta["img_pages"])



                rng = meta["img_pages_ranges"]



                ist = meta["img_stats"]



                img_tag = (f"  [图片转文字 {n} 页: 第 {rng} 页 | "



                           f"图片页字数 {ist['words']} / 中文 {ist['fe']} / "



                           f"非中文 {ist['nc']}]")



            print(f"[OK] {os.path.basename(f)} -> 字数 {s['words']} | "



                  f"中文字符和朝鲜语单词 {s['fe']} | 非中文单词 {s['nc']} | "



                  f"字符数(不计空格) {s['chars']}"



                  + (f"  [CAD OLE={meta.get('ole_count',0)}]" if meta.get('ole_count') else "")



                  + img_tag)



        except Exception as e:



            print(f"[ERR] {os.path.basename(f)} 处理失败: {e}")







    print("\n===== SUMMARY =====")



    tot = {"words": 0, "fe": 0, "nc": 0, "chars": 0}



    for r in results:



        s = r["stats"]



        for k in tot:



            tot[k] += s[k]



        print(f"{os.path.basename(r['src'])}: 字数={s['words']} "



              f"中文字符和朝鲜语单词={s['fe']} 非中文单词={s['nc']} 字符数={s['chars']}")



    print(f"合计: 字数={tot['words']} 中文字符和朝鲜语单词={tot['fe']} "



          f"非中文单词={tot['nc']} 字符数={tot['chars']}")







    png_path = None



    if args.summary_png or args.merge:



        png_path = os.path.join(out_dir, "BATCH_wordcount_summary.png")



        try:



            render_summary_png(results, png_path)



            print("summary png:", png_path)



        except Exception as e:



            print(f"[WARN] 汇总 PNG 生成失败: {e}")



            png_path = None







    if args.merge:



        merged = os.path.join(out_dir, "WordCount_字数统计汇总.docx")



        try:



            build_merged_docx(results, merged, png_path)



            print("merged docx:", merged)



        except Exception as e:



            print(f"[ERR] 合并文档生成失败: {e}")







    if args.send:



        prepare_send(out_dir, results, png_path, args.to, args.subject, args.body, merged if args.merge else None)







    print(f"\n产物目录: {out_dir}")











# ═══════════════════════════════════════════════════════════════════════════
# DOCX 文档比较（仿 Word「审阅 → 比较」）— 纯标准库实现
# 仅使用 Python 标准库：zipfile / xml.etree.ElementTree / json / re / difflib 等
# 不依赖 python-docx / lxml（两者在 Android Chaquopy 上触发 AssetFinder 崩溃）
# 输入：原文档路径、修订文档路径、输出路径、选项 JSON
# 输出：带 w:ins/w:del 修订标记的 .docx + 修改句字数统计
# ═══════════════════════════════════════════════════════════════════════════

# ── WordprocessingML 命名空间（Clark 标记法，供 xml.etree.ElementTree 使用）──
_W = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_WP = "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
_MC = "http://schemas.openxmlformats.org/markup-compatibility/2006"
_CT = "http://schemas.openxmlformats.org/package/2006/content-types"
_REL = "http://schemas.openxmlformats.org/package/2006/relationships"

def _q(tag):
    """返回带命名空间前缀的完整标签名，如 _q('w:p') → '{http://...}p'."""
    prefix, local = tag.split(":")
    ns = {"w": _W, "r": _R, "wp": _WP, "mc": _MC, "ct": _CT, "rel": _REL}[prefix]
    return f"{{{ns}}}{local}"


import xml.etree.ElementTree as ET


def _read_docx_paragraphs(docx_path):
    """用 zipfile+ET 从 DOCX 中提取所有段落文本，返回 [text_str]。
    
    只读取 body 中的顶层 w:p（表格单元格内的段落归入表格文本）。
    """
    import zipfile
    paragraphs = []
    with zipfile.ZipFile(docx_path, "r") as z:
        data = z.read("word/document.xml")
    root = ET.fromstring(data)
    body = root.find(_q("w:body"))
    if body is None:
        return paragraphs
    for child in body:
        if child.tag == _q("w:p"):
            texts = []
            for t_elem in child.iter(_q("w:t")):
                if t_elem.text:
                    texts.append(t_elem.text)
            paragraphs.append("".join(texts))
        elif child.tag == _q("w:tbl"):
            # 表格：汇总所有单元格文本为一个字符串
            cell_texts = []
            for tr in child.iter(_q("w:tr")):
                for tc in tr.iter(_q("w:tc")):
                    for p in tc.iter(_q("w:p")):
                        t_parts = []
                        for t in p.iter(_q("w:t")):
                            if t.text:
                                t_parts.append(t.text)
                        cell_texts.append("".join(t_parts))
            paragraphs.append("\n".join(cell_texts))
    return paragraphs


def _read_docx_paragraphs_with_types(docx_path):
    """返回 [(type, text), ...] 其中 type 为 'p' 或 'tbl'。"""
    import zipfile
    blocks = []
    with zipfile.ZipFile(docx_path, "r") as z:
        data = z.read("word/document.xml")
    root = ET.fromstring(data)
    body = root.find(_q("w:body"))
    if body is None:
        return blocks
    for child in body:
        if child.tag == _q("w:p"):
            texts = []
            for t_elem in child.iter(_q("w:t")):
                if t_elem.text:
                    texts.append(t_elem.text)
            blocks.append(("p", "".join(texts)))
        elif child.tag == _q("w:tbl"):
            cell_texts = []
            for tr in child.iter(_q("w:tr")):
                for tc in tr.iter(_q("w:tc")):
                    for p in tc.iter(_q("w:p")):
                        t_parts = []
                        for t in p.iter(_q("w:t")):
                            if t.text:
                                t_parts.append(t.text)
                        cell_texts.append("".join(t_parts))
            blocks.append(("tbl", "\n".join(cell_texts)))
    return blocks


def _extract_extra_text(docx_path, kind):
    """从 DOCX 中提取附加区域文本（页眉页脚/脚注/文本框/域）。"""
    import zipfile
    parts = []
    try:
        with zipfile.ZipFile(docx_path, "r") as z:
            data = z.read("word/document.xml")
        root = ET.fromstring(data)
        if kind == "header_footer":
            # 遍历 sectPr 中的 headerReference / footerReference 对应文件
            for hf_ref in root.iter(_q("w:headerReference")):
                r_id = hf_ref.get(_q("r:id"), "")
                if r_id:
                    parts.append(f"[header:{r_id}]")
            for hf_ref in root.iter(_q("w:footerReference")):
                r_id = hf_ref.get(_q("r:id"), "")
                if r_id:
                    parts.append(f"[footer:{r_id}]")
        elif kind == "footnote":
            for fn in root.iter(_q("w:footnote")):
                txts = [t.text or "" for t in fn.iter(_q("w:t"))]
                txt = "".join(txts).strip()
                if txt:
                    parts.append(txt)
        elif kind == "textbox":
            for txbx in root.iter(_q("w:txbxContent")):
                txts = [t.text or "" for t in txbx.iter(_q("w:t"))]
                txt = "".join(txts).strip()
                if txt:
                    parts.append(txt)
        elif kind == "field":
            for fld in root.iter(_q("w:fldSimple")):
                txts = [t.text or "" for t in fld.iter(_q("w:t"))]
                txt = "".join(txts).strip()
                if txt:
                    parts.append(txt)
    except Exception:
        pass
    return "\n".join(parts)


def _make_run_element(text, kind, author, date, rid):
    """构建一个 w:r 元素（或包裹在 w:ins/w:del 中）。
    
    kind: 'ins' | 'del' | None(普通文本)
    返回 Element。
    """
    if kind == "ins":
        ins = ET.Element(_q("w:ins"))
        ins.set(_q("w:id"), str(rid))
        ins.set(_q("w:author"), author)
        ins.set(_q("w:date"), date)
        r = ET.SubElement(ins, _q("w:r"))
        rpr = ET.SubElement(r, _q("w:rPr"))
        color = ET.SubElement(rpr, _q("w:color"))
        color.set(_q("w:val"), "2E74B5")  # 插入蓝
        u = ET.SubElement(rpr, _q("w:u"))
        u.set(_q("w:val"), "single")
        t = ET.SubElement(r, _q("w:t"))
        t.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
        t.text = text
        return ins
    elif kind == "del":
        dele = ET.Element(_q("w:del"))
        dele.set(_q("w:id"), str(rid))
        dele.set(_q("w:author"), author)
        dele.set(_q("w:date"), date)
        r = ET.SubElement(dele, _q("w:r"))
        rpr = ET.SubElement(r, _q("w:rPr"))
        strike = ET.SubElement(rpr, _q("w:strike"))
        color = ET.SubElement(rpr, _q("w:color"))
        color.set(_q("w:val"), "C00000")  # 删除红
        t = ET.SubElement(r, _q("w:t"))
        t.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
        t.text = text
        return dele
    else:
        r = ET.Element(_q("w:r"))
        t = ET.SubElement(r, _q("w:t"))
        t.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
        t.text = text
        return r


def _build_diff_paragraph_xml(text_o, text_r, level, author, date,
                               rid_counter, case_sensitive, ignore_ws):
    """对两个段落文本做字符/词级 diff，返回 (p_element, del_ranges)。"""
    import difflib
    p = ET.Element(_q("w:p"))
    toks_o = _normalize_tokens(_tokenize_text(text_o, level), case_sensitive, ignore_ws)
    toks_r = _normalize_tokens(_tokenize_text(text_r, level), case_sensitive, ignore_ws)
    sm = difflib.SequenceMatcher(a=[t[1] for t in toks_o], b=[t[1] for t in toks_r])
    del_ranges = []
    rid_val = rid_counter()
    for tag, i1, i2, j1, j2 in sm.get_opcodes():
        if tag == "equal":
            seg = "".join(t[0] for t in toks_o[i1:i2])
            if seg:
                p.append(_make_run_element(seg, None, author, date, rid_val))
        elif tag == "delete":
            seg = "".join(t[0] for t in toks_o[i1:i2])
            if seg:
                p.append(_make_run_element(seg, "del", author, date, rid_val))
                s = toks_o[i1][2]
                e = toks_o[i2 - 1][3]
                del_ranges.append((s, e))
        elif tag == "insert":
            seg = "".join(t[0] for t in toks_r[j1:j2])
            if seg:
                p.append(_make_run_element(seg, "ins", author, date, rid_val))
        elif tag == "replace":
            dseg = "".join(t[0] for t in toks_o[i1:i2])
            iseg = "".join(t[0] for t in toks_r[j1:j2])
            if dseg:
                p.append(_make_run_element(dseg, "del", author, date, rid_val))
                s = toks_o[i1][2]
                e = toks_o[i2 - 1][3]
                del_ranges.append((s, e))
            if iseg:
                p.append(_make_run_element(iseg, "ins", author, date, rid_val))
    return p, del_ranges


def _build_deleted_paragraph_xml(text, author, date, rid_counter):
    """构建整段删除的 w:p 元素。"""
    p = ET.Element(_q("w:p"))
    if text:
        p.append(_make_run_element(text, "del", author, date, rid_counter()))
    return p


def _build_inserted_paragraph_xml(text, author, date, rid_counter):
    """构建整段插入的 w:p 元素。"""
    p = ET.Element(_q("w:p"))
    if text:
        p.append(_make_run_element(text, "ins", author, date, rid_counter()))
    return p


def _build_note_paragraph_xml(label):
    """构建一个标签段落（如【页眉/页脚变更】）。"""
    p = ET.Element(_q("w:p"))
    r = ET.SubElement(p, _q("w:r"))
    rpr = ET.SubElement(r, _q("w:rPr"))
    b = ET.SubElement(rpr, _q("w:b"))
    t = ET.SubElement(r, _q("w:t"))
    t.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
    t.text = label
    return p


def _build_plain_paragraph_xml(text):
    """构建一个纯文本段落（无修订标记）。"""
    p = ET.Element(_q("w:p"))
    if text:
        r = ET.SubElement(p, _q("w:r"))
        t = ET.SubElement(r, _q("w:t"))
        t.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
        t.text = text
    return p


def _tokenize_text(text, level):
    """返回 list of (orig_text, norm_text, start, end)。"""
    if level == "char":
        return [(c, c, i, i + 1) for i, c in enumerate(text)]
    tokens = []
    i = 0
    n = len(text)
    while i < n:
        ch = text[i]
        if ("\u4e00" <= ch <= "\u9fff") or ("\u3400" <= ch <= "\u4dbf"):
            tokens.append((ch, ch, i, i + 1))
            i += 1
        elif ch.isalnum() or ch == "_":
            j = i
            while j < n and (text[j].isalnum() or text[j] == "_"):
                j += 1
            tokens.append((text[i:j], text[i:j], i, j))
            i = j
        else:
            tokens.append((ch, ch, i, i + 1))
            i += 1
    return tokens


def _normalize_tokens(tokens, case_sensitive, ignore_ws):
    out = []
    for orig, _, s, e in tokens:
        norm = orig
        if not case_sensitive:
            norm = norm.lower()
        if ignore_ws:
            norm = re.sub(r"\s+", "", norm)
        out.append((orig, norm, s, e))
    return out


def _split_sentences(text):
    parts = re.split(r"([。！？；\n\r])", text)
    res = []
    pos = 0
    i = 0
    n = len(parts)
    while i < n:
        seg = parts[i]
        if i + 1 < n:
            full = seg + parts[i + 1]
            if full.strip():
                res.append((pos, pos + len(full), full))
            pos += len(full)
            i += 2
        else:
            if seg.strip():
                res.append((pos, pos + len(seg), seg))
            break
    return res


def _count_modified_sentences(text, ranges):
    if not text.strip() or not ranges:
        return 0
    total = 0
    for (s, e, sent) in _split_sentences(text):
        for (rs, re_) in ranges:
            if e > rs and s < re_:
                total += len(re.sub(r"\s", "", sent))
                break
    return total


def _count_text_chars(text):
    return len(re.sub(r"\s", "", text or ""))


def _create_output_docx(orig_path, out_path, body_elements):
    """基于原文档创建输出 DOCX，替换 body 内容为 body_elements。
    
    body_elements: list of ET.Element（w:p 或其他 body 子元素）
    保留原文档的所有样式/rels/内容类型等，只替换 document.xml 的 body。
    """
    import zipfile

    # 注册命名空间前缀（让输出 XML 可读且 Word 兼容）
    ET.register_namespace("w", _W)
    ET.register_namespace("r", _R)
    ET.register_namespace("wp", _WP)
    ET.register_namespace("mc", _MC)
    ET.register_namespace("v", "urn:schemas-microsoft-com:vml")
    ET.register_namespace("o", "urn:schemas-microsoft-com:office:office")
    ET.register_namespace("m", "http://schemas.openxmlformats.org/officeDocument/2006/math")

    # 读取原文档
    with zipfile.ZipFile(orig_path, "r") as zin:
        # 解析 document.xml
        doc_data = zin.read("word/document.xml")
        root = ET.fromstring(doc_data)

        # 替换 body 内容
        body = root.find(_q("w:body"))
        if body is None:
            body = ET.SubElement(root, _q("w:body"))

        # 移除现有子元素（保留 sectPr 如果存在）
        sect_pr = None
        existing = list(body)
        for child in existing:
            if child.tag == _q("w:sectPr"):
                sect_pr = child
            body.remove(child)

        # 添加新的 body 元素
        for elem in body_elements:
            body.append(elem)

        # 恢复 sectPr（必须在最后）
        if sect_pr is not None:
            body.append(sect_pr)

        # 添加 trackRevisions 到 settings.xml（如果存在）
        settings_data = None
        new_settings_data = None
        try:
            settings_data = zin.read("word/settings.xml")
            sroot = ET.fromstring(settings_data)
            # 移除已有的 trackRevisions
            for tr in sroot.findall(_q("w:trackRevisions")):
                sroot.remove(tr)
            # 添加新的
            tr = ET.SubElement(sroot, _q("w:trackRevisions"))
            new_settings_data = ET.tostring(sroot, encoding="unicode", xml_declaration=True)
            # ET.tostring 返回 str，需要编码为 bytes
            if isinstance(new_settings_data, str):
                new_settings_data = new_settings_data.encode("utf-8")
        except Exception:
            pass

        # 写出输出文件
        new_doc_data = ET.tostring(root, encoding="unicode", xml_declaration=True)
        if isinstance(new_doc_data, str):
            new_doc_data = new_doc_data.encode("utf-8")

        # 构建输出 ZIP（复制原文档 + 替换修改过的文件）
        with zipfile.ZipFile(out_path, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                if item.filename == "word/document.xml":
                    zout.writestr(item, new_doc_data)
                elif item.filename == "word/settings.xml" and new_settings_data is not None:
                    zout.writestr(item, new_settings_data)
                else:
                    zout.writestr(item, zin.read(item.filename))


def compare_docx(orig_path, rev_path, out_path, opts_json):
    """比较两份 DOCX，生成修订标记 .docx 并统计修改句字数。

    纯标准库实现（zipfile + xml.etree.ElementTree），不依赖 python-docx/lxml。
    返回 JSON 字符串.
    """
    import json, re, datetime, difflib, sys, traceback, os

    try:
        # ── 前置校验：文件存在性 ──
        if not os.path.isfile(orig_path):
            return json.dumps({"ok": False, "error": f"原文档不存在: {orig_path}"}, ensure_ascii=False)
        if not os.path.isfile(rev_path):
            return json.dumps({"ok": False, "error": f"修订文档不存在: {rev_path}"}, ensure_ascii=False)

        # ── 解析选项 ──
        try:
            opts = json.loads(opts_json) if opts_json else {}
        except Exception:
            opts = {}
        level = opts.get("level", "word")
        case_sensitive = opts.get("case", True)
        ignore_ws = opts.get("whitespace", False)
        use_table = opts.get("table", True)
        use_hf = opts.get("header_footer", True)
        use_fn = opts.get("footnote", True)
        use_tb = opts.get("textbox", True)
        use_field = opts.get("field", True)

        author = "WordCount"
        date = datetime.datetime.now().strftime("%Y-%m-%dT%H:%M:%SZ")

        # ── 读取两份文档的块（段落 + 表格）──
        try:
            blocks_o = _read_docx_paragraphs_with_types(orig_path)
        except Exception as e:
            return json.dumps({"ok": False, "error": f"读取原文档失败: {type(e).__name__}: {str(e)[:300]}"}, ensure_ascii=False)
        try:
            blocks_r = _read_docx_paragraphs_with_types(rev_path)
        except Exception as e:
            return json.dumps({"ok": False, "error": f"读取修订文档失败: {type(e).__name__}: {str(e)[:300]}"}, ensure_ascii=False)

        # ── difflib 段落级比较 ──
        sm = difflib.SequenceMatcher(
            a=[b[1] for b in blocks_o],
            b=[b[1] for b in blocks_r]
        )
        opcodes = sm.get_opcodes()

        # 合并相邻 delete+insert 为 replace
        merged = []
        ki = 0
        while ki < len(opcodes):
            tag, i1, i2, j1, j2 = opcodes[ki]
            if tag == "delete" and ki + 1 < len(opcodes):
                t2, a1, a2, b1, b2 = opcodes[ki + 1]
                if t2 == "insert" and (i2 - i1) == (b2 - b1):
                    merged.append(("replace", i1, i2, j1, b2))
                    ki += 2
                    continue
            if tag == "insert" and ki + 1 < len(opcodes):
                t2, a1, a2, b1, b2 = opcodes[ki + 1]
                if t2 == "delete" and (a2 - a1) == (j2 - j1):
                    merged.append(("replace", a1, a2, j1, j2))
                    ki += 2
                    continue
            merged.append((tag, i1, i2, j1, j2))
            ki += 1

        # ── 构建 output body 元素 ──
        rid_seq = [0]

        def rid_counter():
            rid_seq[0] += 1
            return rid_seq[0]

        body_elements = []
        modified_chars = 0
        ins_count = 0
        del_count = 0
        rep_count = 0

        for tag, i1, i2, j1, j2 in merged:
            if tag == "equal":
                for k in range(i1, i2):
                    el = blocks_o[k]
                    if el[0] == "p":
                        body_elements.append(_build_plain_paragraph_xml(el[1]))
                    else:
                        body_elements.append(_build_note_paragraph_xml(f"[原文表格] {el[1][:200]}"))
            elif tag == "delete":
                for k in range(i1, i2):
                    el = blocks_o[k]
                    if el[0] == "p":
                        body_elements.append(_build_deleted_paragraph_xml(el[1], author, date, rid_counter))
                        modified_chars += _count_modified_sentences(el[1], [(0, len(el[1]))])
                    else:
                        body_elements.append(_build_note_paragraph_xml(f"[已删除表格] {el[1][:200]}"))
                        modified_chars += _count_text_chars(el[1])
                    del_count += 1
            elif tag == "insert":
                for k in range(j1, j2):
                    el = blocks_r[k]
                    if el[0] == "p":
                        body_elements.append(_build_inserted_paragraph_xml(el[1], author, date, rid_counter))
                    else:
                        body_elements.append(_build_note_paragraph_xml(f"[新增表格] {el[1][:200]}"))
                    ins_count += 1
            elif tag == "replace":
                single = (i2 - i1 == 1 and j2 - j1 == 1)
                if single:
                    bo = blocks_o[i1]
                    br = blocks_r[j1]
                    if bo[0] == "p" and br[0] == "p":
                        p, ranges = _build_diff_paragraph_xml(
                            bo[1], br[1], level, author, date,
                            rid_counter, case_sensitive, ignore_ws)
                        body_elements.append(p)
                        modified_chars += _count_modified_sentences(bo[1], ranges)
                        rep_count += 1
                        continue
                    if bo[0] == "tbl" and br[0] == "tbl" and use_table:
                        # 表格 diff：降级为占位段落
                        body_elements.append(_build_note_paragraph_xml(f"[原表格] {bo[1][:200]}"))
                        body_elements.append(_build_note_paragraph_xml(f"[修订表格] {br[1][:200]}"))
                        modified_chars += _count_text_chars(bo[1])
                        rep_count += 1
                        continue
                # 多块替换或类型不匹配
                for k in range(i1, i2):
                    el = blocks_o[k]
                    if el[0] == "p":
                        body_elements.append(_build_deleted_paragraph_xml(el[1], author, date, rid_counter))
                        modified_chars += _count_modified_sentences(el[1], [(0, len(el[1]))])
                    else:
                        body_elements.append(_build_note_paragraph_xml(f"[已删除表格] {el[1][:200]}"))
                        modified_chars += _count_text_chars(el[1])
                    del_count += 1
                for k in range(j1, j2):
                    el = blocks_r[k]
                    if el[0] == "p":
                        body_elements.append(_build_inserted_paragraph_xml(el[1], author, date, rid_counter))
                    else:
                        body_elements.append(_build_note_paragraph_xml(f"[新增表格] {el[1][:200]}"))
                    ins_count += 1

        # ── 附加区域 ──
        extra_kinds = []
        if use_hf:
            extra_kinds.append(("header_footer", "【页眉/页脚变更】"))
        if use_fn:
            extra_kinds.append(("footnote", "【脚注/尾注变更】"))
        if use_tb:
            extra_kinds.append(("textbox", "【文本框变更】"))
        if use_field:
            extra_kinds.append(("field", "【域变更】"))
        for kind, label in extra_kinds:
            to = _extract_extra_text(orig_path, kind)
            tr = _extract_extra_text(rev_path, kind)
            if to != tr:
                body_elements.append(_build_note_paragraph_xml(label))
                p, _ = _build_diff_paragraph_xml(
                    to, tr, level, author, date,
                    rid_counter, case_sensitive, ignore_ws)
                body_elements.append(p)

        # ── 写出输出 DOCX ──
        try:
            _create_output_docx(orig_path, out_path, body_elements)
        except Exception as e:
            return json.dumps({"ok": False, "error": f"写出结果文档失败: {type(e).__name__}: {str(e)[:300]}"}, ensure_ascii=False)

        return json.dumps({
            "ok": True,
            "out_path": out_path,
            "insertions": ins_count,
            "deletions": del_count,
            "replacements": rep_count,
            "modified_sentence_chars": modified_chars,
        }, ensure_ascii=False)

    except Exception as e:
        err_msg = f"{type(e).__name__}: {str(e)[:400]}"
        tb_text = traceback.format_exc()[:1200]
        sys.stderr.write(f"COMPARE_ERROR: {err_msg}\n{tb_text}\n")
        return json.dumps({"ok": False, "error": err_msg, "trace": tb_text}, ensure_ascii=False)


if __name__ == "__main__":



    main()



